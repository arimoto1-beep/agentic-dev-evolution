"""投稿に使う 1 枚絵(サムネイル)。

画像そのものは LibreOffice が要るため slow に分け、ここでは主に
「何を出すか」(題と副題の取り出し)と「どう置くか」(文字の大きさと行数)を確認する。
"""

import os

import pytest
from pptx import Presentation
from PIL import Image

from note2slides import convert_file
from note2slides.renderer import Renderer, render_deck
from note2slides.slide_images import ImageOptions
from note2slides.soffice import find_soffice
from note2slides.style import Style, get_theme
from note2slides.thumbnail import Thumbnail, ThumbnailError, export_thumbnail, from_source
from note2slides.thumbnail_cli import main

soffice = find_soffice()
requires_soffice = pytest.mark.skipif(
    soffice is None, reason="LibreOffice(soffice)が見つかりません"
)

SCENARIO = """---
type: scenario
title: 教材シナリオから動画を作る
---

## 教材シナリオから動画を作る

### 設定

- レイアウト: 表紙

### 画面

note2slides / サンプル

### ナレーション

はじめます。

## つぎの画面

### 画面

- 項目
"""

ARTICLE = """---
title: eラーニング動画のつくり方
subtitle: 記事から作る
---

## はじめに

本文です。
"""


def write(tmp_path, name, text):
    path = tmp_path / name
    path.write_text(text, encoding="utf-8")
    return str(path)


def render(tmp_path, thumbnail: Thumbnail, theme: str = "light"):
    out = str(tmp_path / "thumb.pptx")
    render_deck(thumbnail.to_deck(), out, style=Style(theme=get_theme(theme)))
    return Presentation(out).slides[0]


def texts(slide):
    return [
        shape.text_frame.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    ]


class TestWhatIsShown:
    def test_title_subtitle_and_label_are_drawn(self, tmp_path):
        slide = render(tmp_path, Thumbnail("教材シナリオから動画を作る", "副題です", "第23回"))

        drawn = " ".join(texts(slide))
        assert "副題です" in drawn
        assert "第23回" in drawn
        assert "教材シナリオ" in drawn

    def test_only_one_slide_is_made(self, tmp_path):
        out = str(tmp_path / "thumb.pptx")
        render_deck(Thumbnail("題").to_deck(), out)

        assert len(Presentation(out).slides) == 1

    def test_there_is_no_page_number(self, tmp_path):
        """サムネイルは資料ではないので、フッタは出さない。"""
        slide = render(tmp_path, Thumbnail("題", "副題"))

        assert not any("1 / 1" in text for text in texts(slide))


class TestHowItIsPlaced:
    def title_lines(self, title: str):
        renderer = Renderer(Style(theme=get_theme("light")))
        return renderer._thumbnail_title(title)

    def test_a_long_title_uses_smaller_letters(self):
        big, _ = self.title_lines("短い題")
        small, _ = self.title_lines("生成AIを前提とした開発環境をゼロから作り直すための実践ガイド")

        assert small < big

    def test_a_long_title_is_not_split_inside_a_word(self):
        _, lines = self.title_lines("生成AIを前提とした開発環境をゼロから作り直すための実践ガイド")

        assert len(lines) == 2
        assert lines[1].startswith("ゼロから")

    def test_titles_stay_within_three_lines(self):
        long_title = (
            "開発環境そのものを改善しながらソフトウェアを作るという実験について、"
            "その進め方と観察結果をまとめました"
        )

        _, lines = self.title_lines(long_title)

        assert 1 <= len(lines) <= 3
        assert "".join(lines) == long_title

    def test_a_short_title_uses_the_biggest_letters(self):
        size, lines = self.title_lines("教材動画のつくり方")

        assert lines == ["教材動画のつくり方"]
        assert size == 68


class TestTakingTheTitleFromTheInput:
    def test_scenario_gives_its_cover(self, tmp_path):
        path = write(tmp_path, "lesson.md", SCENARIO)

        thumbnail = from_source(path)

        assert thumbnail.title == "教材シナリオから動画を作る"
        assert thumbnail.subtitle == "note2slides / サンプル"

    def test_article_gives_its_title_and_subtitle(self, tmp_path):
        path = write(tmp_path, "article.md", ARTICLE)

        thumbnail = from_source(path)

        assert thumbnail.title == "eラーニング動画のつくり方"
        assert thumbnail.subtitle == "記事から作る"

    def test_pptx_gives_the_cover_without_the_footer(self, tmp_path):
        md = write(tmp_path, "article.md", ARTICLE)
        pptx_path = str(tmp_path / "article.pptx")
        convert_file(md, pptx_path)

        thumbnail = from_source(pptx_path)

        assert thumbnail.title == "eラーニング動画のつくり方"
        assert "/" not in thumbnail.subtitle  # ページ番号(1 / 2)を拾わない
        assert thumbnail.subtitle == "記事から作る"

    def test_a_wrapped_title_becomes_one_line(self, tmp_path):
        """表紙で 2 行に折り返した題も、取り出すときは 1 行に戻す。"""
        long_title = "生成AIを前提とした開発環境をゼロから作り直すための実践ガイド"
        md = write(tmp_path, "a.md", f"---\ntitle: {long_title}\n---\n\n## 節\n\n本文。\n")
        pptx_path = str(tmp_path / "a.pptx")
        convert_file(md, pptx_path)

        assert from_source(pptx_path).title == long_title

    def test_unsupported_input_is_reported(self, tmp_path):
        path = write(tmp_path, "note.txt", "本文")

        with pytest.raises(ThumbnailError):
            from_source(path)

    def test_missing_input_is_reported(self, tmp_path):
        with pytest.raises(ThumbnailError):
            from_source(str(tmp_path / "no-such-file.md"))


class TestExport:
    def test_an_empty_title_is_reported(self, tmp_path):
        with pytest.raises(ThumbnailError):
            export_thumbnail(Thumbnail("   "), str(tmp_path / "out.png"))

    def test_an_existing_file_is_not_overwritten(self, tmp_path):
        out = tmp_path / "out.png"
        out.write_bytes(b"")

        with pytest.raises(ThumbnailError):
            export_thumbnail(Thumbnail("題"), str(out))

    @pytest.mark.slow
    @requires_soffice
    def test_the_image_has_the_size_youtube_asks_for(self, tmp_path):
        out = str(tmp_path / "thumbnail.png")

        result = export_thumbnail(
            Thumbnail("教材シナリオから動画を作る", "note2slides", "第23回"),
            out,
            soffice_path=soffice,
        )

        assert (result.width, result.height) == (1280, 720)
        with Image.open(out) as image:
            assert image.size == (1280, 720)
            assert image.mode == "RGB"
            # 地が塗られていること(白紙のままではないこと)。
            assert image.getpixel((10, 10)) != (255, 255, 255)


class TestCommand:
    def test_it_asks_for_a_title_when_there_is_no_input(self, capsys):
        assert main([]) == 2

    def test_an_empty_title_is_refused(self, tmp_path, capsys):
        path = write(tmp_path, "a.md", "---\ntitle: 題\n---\n\n## 節\n\n本文。\n")

        assert main([path, "--title", "   ", "-o", str(tmp_path / "t.png")]) == 2

    @pytest.mark.slow
    @requires_soffice
    def test_it_writes_the_image_next_to_the_input(self, tmp_path):
        path = write(tmp_path, "lesson.md", SCENARIO)

        assert main([path, "--quiet", "--soffice", soffice]) == 0

        assert os.path.isfile(str(tmp_path / "lesson_thumbnail.png"))


def test_image_options_for_the_default_size():
    """既定の大きさが 16:9 で、動画側の決まり(偶数)も満たすこと。"""
    options = ImageOptions(width=1280)

    options.validate()
    assert options.size() == (1280, 720)
