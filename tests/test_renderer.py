"""生成した .pptx を読み直して構造を確認する。"""

import re

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from note2slides import convert_file
from note2slides.markdown_reader import parse_article
from note2slides.planner import PlannerOptions, plan_deck
from note2slides.renderer import render_deck
from note2slides.style import SLIDE_HEIGHT_EMU, SLIDE_WIDTH_EMU

ARTICLE = """---
title: テスト記事
subtitle: 副題
---

# テスト記事

## 概要

一文目です。二文目です。

- 箇条書き **強調** あり
  - 入れ子

## 手順

1. 最初
2. 次

> 引用です。

```python
x = 1
```

| 列A | 列B |
| --- | --- |
| 1 | 2 |
"""


@pytest.fixture
def deck_and_pptx(tmp_path):
    md = tmp_path / "article.md"
    md.write_text(ARTICLE, encoding="utf-8")
    out = tmp_path / "article.pptx"
    deck = convert_file(str(md), str(out))
    return deck, Presentation(str(out))


def test_slide_size_is_16_9(deck_and_pptx):
    _, prs = deck_and_pptx
    assert prs.slide_width == SLIDE_WIDTH_EMU
    assert prs.slide_height == SLIDE_HEIGHT_EMU
    assert round(prs.slide_width / prs.slide_height, 3) == round(16 / 9, 3)


def test_slide_count_matches_plan(deck_and_pptx):
    deck, prs = deck_and_pptx
    assert len(prs.slides) == len(deck.slides)
    assert len(prs.slides) > 1


def test_titles_are_real_title_placeholders(deck_and_pptx):
    deck, prs = deck_and_pptx
    rendered = [s.shapes.title.text_frame.text for s in prs.slides if s.shapes.title is not None]
    assert "テスト記事" in rendered
    assert "概要" in rendered
    assert "手順" in rendered


def test_notes_carry_source_text(deck_and_pptx):
    _, prs = deck_and_pptx
    notes = [s.notes_slide.notes_text_frame.text for s in prs.slides if s.has_notes_slide]
    assert any("一文目です。二文目です。" in n for n in notes)


def test_table_is_rendered_as_table(deck_and_pptx):
    _, prs = deck_and_pptx
    tables = [sh.table for s in prs.slides for sh in s.shapes if sh.has_table]
    assert len(tables) == 1
    assert tables[0].cell(0, 0).text_frame.text == "列A"
    assert tables[0].cell(1, 1).text_frame.text == "2"


def test_bullet_formatting_is_written(deck_and_pptx):
    _, prs = deck_and_pptx
    xml = "".join(s.shapes._spTree.xml for s in prs.slides)
    assert "buChar" in xml  # 箇条書き記号
    assert "buAutoNum" in xml  # 番号付きリスト
    assert "buNone" in xml  # タイトルやコードなど記号なし


def test_east_asian_font_is_specified(deck_and_pptx):
    _, prs = deck_and_pptx
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    rPr = run._r.find(qn("a:rPr"))
                    assert rPr is not None and rPr.find(qn("a:ea")) is not None


def _normalize(text: str) -> str:
    """比較用に、空白と Markdown の記号を落とす。"""
    text = text.replace("（続き）", "")
    return re.sub(r"[\s#*`>|_~\-]+", "", text)


def _text_units(prs):
    """スライド上の文字列を、段落・セル単位で取り出す。"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    yield "".join(run.text for run in paragraph.runs)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame.text


def test_generated_text_comes_from_the_article(deck_and_pptx):
    """スライド上の文字列が、記事に無い内容になっていないことを確認する。"""
    _, prs = deck_and_pptx
    source = _normalize(ARTICLE)
    for text in _text_units(prs):
        normalized = _normalize(text)
        if not normalized:
            continue
        assert normalized in source, text


def test_convert_file_returns_deck(tmp_path):
    md = tmp_path / "a.md"
    md.write_text("## 節\n\n本文。\n", encoding="utf-8")
    out = tmp_path / "a.pptx"
    deck = convert_file(str(md), str(out))
    assert out.exists()
    assert deck.slides


def test_image_is_embedded(tmp_path):
    from PIL import Image as PILImage  # python-pptx の依存として利用できる

    png = tmp_path / "figure.png"
    PILImage.new("RGB", (640, 360), (220, 220, 220)).save(png)
    md_text = "## 図\n\n![キャプション](figure.png)\n"
    article = parse_article(md_text, source_path=str(tmp_path / "a.md"))
    deck = plan_deck(article, options=PlannerOptions(title_slide=False))
    out = tmp_path / "a.pptx"
    render_deck(deck, str(out))

    prs = Presentation(str(out))
    pictures = [sh for s in prs.slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    assert not deck.warnings
    picture = pictures[0]
    assert picture.width <= SLIDE_WIDTH_EMU
    assert picture.top + picture.height <= SLIDE_HEIGHT_EMU


def test_line_break_inside_a_paragraph_is_a_real_break(tmp_path):
    """記事の中の改行は、段落を分けずに行だけを変える(`<a:br/>`)。

    `<a:t>` の中に改行文字を置いても描画側は空白として扱うため、note の `<br>`
    で書かれた図(文章 ↓ トークン …)が 1 行につながってしまう。
    """
    article = parse_article("## 図\n\n文章  \n↓  \nトークン\n")
    deck = plan_deck(article, options=PlannerOptions(title_slide=False))
    out = tmp_path / "br.pptx"
    render_deck(deck, str(out))

    prs = Presentation(str(out))
    paragraphs = [
        p
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for p in shape.text_frame.paragraphs
        if "".join(r.text for r in p.runs).startswith("文章")
    ]
    assert len(paragraphs) == 1
    body = paragraphs[0]
    assert len(body._p.findall(qn("a:br"))) == 2
    assert [r.text for r in body.runs] == ["文章", "↓", "トークン"]
    # python-pptx は <a:br/> を垂直タブとして読み直す(読み上げ側はここで行を分ける)。
    assert body.text == "文章\v↓\vトークン"
