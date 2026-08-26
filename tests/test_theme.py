"""スライドの見た目(テーマ)と、タイトルの折り返し。

見た目そのものは目で見て決めるものなので、ここでは「見た目を変えても壊れては
いけないこと」を確認する。地の色が付くこと、飾りの文字が読み上げに混ざらない
こと、従来の見た目(plain)がそのまま残っていること。
"""

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from note2slides.model import (
    KIND_BULLETS,
    KIND_SECTION,
    KIND_TITLE,
    SHAPE_FOOTER,
    Bullet,
    Deck,
    Run,
    Slide,
    parse_shape_name,
)
from note2slides.narration import extract_script
from note2slides.renderer import render_deck
from note2slides.style import THEMES, Style, get_theme, theme_names

LONG_TITLE = "生成AIを前提とした開発環境をゼロから作り直すための実践ガイド"


def make_deck(title: str = "資料の題") -> Deck:
    return Deck(
        title=title,
        slides=[
            Slide(kind=KIND_TITLE, title=title, subtitle="副題"),
            Slide(kind=KIND_SECTION, title="第1章"),
            Slide(
                kind=KIND_BULLETS,
                title="見出し",
                bullets=[Bullet(runs=[Run("1 つ目の項目")]), Bullet(runs=[Run("2 つ目の項目")])],
            ),
        ],
    )


def render(tmp_path, deck: Deck, theme: str = "light"):
    out = tmp_path / f"{theme}.pptx"
    render_deck(deck, str(out), style=Style(theme=get_theme(theme)))
    return Presentation(str(out))


def background_rgb(slide):
    """スライドに塗られた地の色(塗っていなければ None)。"""
    fill = slide.background._element.find(qn("p:bg"))
    if fill is None:
        return None
    color = fill.find(".//" + qn("a:srgbClr"))
    return None if color is None else color.get("val")


def footer_texts(slide):
    return [
        shape.text_frame.text
        for shape in slide.shapes
        if parse_shape_name(shape.name)[0] == SHAPE_FOOTER
    ]


class TestTheme:
    def test_cover_and_section_have_their_own_background(self, tmp_path):
        prs = render(tmp_path, make_deck())

        cover, section, body = prs.slides
        assert background_rgb(cover) is not None
        assert background_rgb(section) is not None
        assert background_rgb(cover) != background_rgb(section)
        # 本文は白のまま(地を塗らない)。
        assert background_rgb(body) is None

    def test_page_number_is_on_body_slides_but_not_on_the_cover(self, tmp_path):
        prs = render(tmp_path, make_deck())

        cover, section, body = prs.slides
        assert footer_texts(cover) == []
        assert "2 / 3" in footer_texts(section)
        assert "3 / 3" in footer_texts(body)
        assert "資料の題" in footer_texts(body)

    def test_long_deck_title_is_shortened_in_the_footer(self, tmp_path):
        prs = render(tmp_path, make_deck(LONG_TITLE * 2))

        footer = footer_texts(prs.slides[2])
        assert any(text.endswith("…") for text in footer)

    def test_footer_is_not_read_as_narration(self, tmp_path):
        """飾りの文字(資料名・ページ番号)を読み上げてしまわないこと。"""
        out = tmp_path / "deck.pptx"
        render_deck(make_deck(), str(out), style=Style(theme=get_theme("light")))

        segment = extract_script(str(out)).segments[2]

        assert "3 / 3" not in segment.text
        assert segment.text == "見出し\n1 つ目の項目\n2 つ目の項目"

    def test_plain_theme_keeps_the_previous_look(self, tmp_path):
        prs = render(tmp_path, make_deck(), theme="plain")

        for slide in prs.slides:
            assert background_rgb(slide) is None
            assert footer_texts(slide) == []

    def test_plain_keeps_the_cover_where_it_was(self, tmp_path):
        """plain の表紙は、装飾を足す前と同じ場所に置く。"""
        (tmp_path / "a").mkdir()
        (tmp_path / "b").mkdir()
        plain = render(tmp_path / "a", make_deck(), theme="plain").slides[0]
        light = render(tmp_path / "b", make_deck()).slides[0]

        title = plain.shapes.title
        assert round((title.top + title.height) / 914400, 2) == 3.8
        # light は下端の帯があるぶん、少し上に置く。
        assert light.shapes.title.top < title.top

    def test_unknown_theme_is_reported(self):
        with pytest.raises(ValueError):
            get_theme("むらさき")

    def test_known_themes_are_listed(self):
        assert set(theme_names()) == set(THEMES)
        assert "light" in theme_names() and "plain" in theme_names()


class TestTitleWrapping:
    def title_lines(self, slide):
        """タイトル図形の中の行(pptx の行内改行は `<a:br/>`)。"""
        paragraph = slide.shapes.title.text_frame.paragraphs[0]
        lines, current = [], ""
        for child in paragraph._p:
            if child.tag == qn("a:r"):
                current += child.text
            elif child.tag == qn("a:br"):
                lines.append(current)
                current = ""
        lines.append(current)
        return lines

    def test_long_cover_title_is_broken_at_a_natural_place(self, tmp_path):
        prs = render(tmp_path, make_deck(LONG_TITLE))

        lines = self.title_lines(prs.slides[0])

        assert len(lines) == 2
        assert "".join(lines) == LONG_TITLE
        # 語の途中(「ゼロか / ら」)ではなく、文節の切れ目で分ける。
        assert lines[1].startswith("ゼロから")

    def test_short_title_stays_on_one_line(self, tmp_path):
        prs = render(tmp_path, make_deck("短い題"))

        assert self.title_lines(prs.slides[0]) == ["短い題"]

    def test_the_title_box_grows_with_the_number_of_lines(self, tmp_path):
        (tmp_path / "a").mkdir()
        (tmp_path / "b").mkdir()
        one = render(tmp_path / "a", make_deck("短い題"))
        two = render(tmp_path / "b", make_deck(LONG_TITLE))

        short_box = one.slides[0].shapes.title
        long_box = two.slides[0].shapes.title
        assert long_box.height > short_box.height
        # 下端(副題との間)は動かさない。
        assert long_box.top + long_box.height == short_box.top + short_box.height


class TestWhatIsWrittenIsWhatAppears:
    """ひな型のレイアウトが、書いた文字を勝手に変えないこと。

    python-pptx の既定のひな型で章扉に使う `Section Header` は、タイトルに
    `cap="all"` を持っている。何も書かないとこれが効き、「試験Runner」が
    「試験RUNNER」として出る。資料の中の文字は正しいままなので、
    スライド画像を目で見るまで気付かない。
    """

    def _section_title_run(self, tmp_path, title):
        deck = Deck(slides=[Slide(kind=KIND_SECTION, title=title)], title="題")
        path = str(tmp_path / "deck.pptx")
        render_deck(deck, path, Style())
        pptx = Presentation(path)
        for shape in pptx.slides[0].shapes:
            if shape.has_text_frame and title in shape.text_frame.text:
                return shape.text_frame.paragraphs[0].runs[0]
        raise AssertionError("章扉の見出しが見つかりません")

    def test_a_section_title_is_not_turned_into_capitals(self, tmp_path):
        run = self._section_title_run(tmp_path, "おまけ: 試験Runnerを作るなら")

        assert run._r.get_or_add_rPr().get("cap") == "none"

    def test_the_layout_this_guards_against_really_asks_for_capitals(self, tmp_path):
        # 守っている相手が実在することの確認。ひな型が変わってこれが通らなくなったら、
        # 上のテストは何も守っていないことになる。
        deck = Deck(slides=[Slide(kind=KIND_SECTION, title="見出し")], title="題")
        path = str(tmp_path / "deck.pptx")
        render_deck(deck, path, Style())

        layout = Presentation(path).slides[0].slide_layout
        assert 'cap="all"' in layout.element.xml
