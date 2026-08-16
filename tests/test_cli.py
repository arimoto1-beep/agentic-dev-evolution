"""資料生成コマンドの入力(ファイル / 公開 note 記事の URL)。

URL からの取り込みは通信を差し替えて確認する(実際の通信は
`test_note_source_live.py`)。
"""

from __future__ import annotations

import json

from pptx import Presentation

from note2slides import cli
from test_note_source import ARTICLE_URL, fake_network, note_payload

BODY = (
    "<p>生成AIを使っていると、トークンという言葉を見かけます。</p>"
    "<h2>AIは文章をそのまま読んでいない</h2>"
    "<p>まず文章を一定の単位に分割します。</p>"
    "<figure><img src='https://a.test/1.png'><figcaption>流れの図</figcaption></figure>"
)


def test_a_markdown_file_still_works(tmp_path):
    article = tmp_path / "a.md"
    article.write_text("## 節\n\n本文です。\n", encoding="utf-8")
    out = tmp_path / "a.pptx"

    assert cli.main([str(article), "-o", str(out), "--quiet"]) == cli.EXIT_OK
    assert out.exists()


def test_a_missing_file_is_reported(tmp_path, capsys):
    assert cli.main([str(tmp_path / "none.md"), "--quiet"]) == cli.EXIT_USAGE
    assert "見つかりません" in capsys.readouterr().err


class TestFromScenario:
    """教材シナリオは、記事と同じコマンドで資料になる(front matter で見分ける)。"""

    SCENARIO = (
        "---\ntype: scenario\ntitle: 教材\n---\n\n"
        "## 見出し\n\n### 画面\n- 一つ目\n- 二つ目\n\n"
        "### ナレーション\n二つの点を説明します。\n"
    )

    def write(self, tmp_path, text=None, name="lesson.md"):
        path = tmp_path / name
        path.write_text(text or self.SCENARIO, encoding="utf-8")
        return path

    def test_a_scenario_becomes_a_presentation(self, tmp_path):
        out = tmp_path / "lesson.pptx"

        assert cli.main([str(self.write(tmp_path)), "-o", str(out), "--quiet"]) == cli.EXIT_OK

        prs = Presentation(str(out))
        assert prs.core_properties.title == "教材"
        assert len(prs.slides) == 1

    def test_the_narration_goes_into_the_presenter_notes(self, tmp_path):
        out = tmp_path / "lesson.pptx"

        cli.main([str(self.write(tmp_path)), "-o", str(out), "--quiet"])

        notes = Presentation(str(out)).slides[0].notes_slide.notes_text_frame.text
        assert notes == "二つの点を説明します。"

    def test_the_screen_is_not_rebuilt_from_the_narration(self, tmp_path):
        """画面に出るのはシナリオの `### 画面` だけ(ナレーションは画面に出ない)。"""
        plan = tmp_path / "plan.json"

        cli.main(
            [
                str(self.write(tmp_path)),
                "-o",
                str(tmp_path / "lesson.pptx"),
                "--dump-plan",
                str(plan),
                "--quiet",
            ]
        )

        slides = json.loads(plan.read_text(encoding="utf-8"))["slides"]
        assert [b["text"] for b in slides[0]["bullets"]] == ["一つ目", "二つ目"]

    def test_a_mistake_in_the_scenario_stops_before_writing(self, tmp_path, capsys):
        path = self.write(tmp_path, "---\ntype: scenario\n---\n\n## 見出し\n本文\n")
        out = tmp_path / "lesson.pptx"

        assert cli.main([str(path), "-o", str(out), "--quiet"]) == cli.EXIT_USAGE

        assert "### 画面" in capsys.readouterr().err
        assert not out.exists()

    def test_article_only_options_are_refused(self, tmp_path, capsys):
        path = self.write(tmp_path)

        code = cli.main([str(path), "-o", str(tmp_path / "a.pptx"), "--no-notes", "--quiet"])

        assert code == cli.EXIT_USAGE
        assert "--no-notes" in capsys.readouterr().err

    def test_asking_for_a_scenario_that_is_an_article_says_so(self, tmp_path, capsys):
        path = self.write(tmp_path, "# 記事\n\n本文です。\n", name="article.md")

        code = cli.main([str(path), "--scenario", "-o", str(tmp_path / "a.pptx"), "--quiet"])

        assert code == cli.EXIT_USAGE
        assert "type: scenario" in capsys.readouterr().err

    def test_an_article_still_goes_through_the_planner(self, tmp_path):
        """シナリオの見分けは front matter だけ。記事の扱いは変わらない。"""
        path = self.write(tmp_path, "## 節\n\n一つ目の文です。二つ目の文です。\n", name="a.md")
        plan = tmp_path / "plan.json"

        cli.main([str(path), "-o", str(tmp_path / "a.pptx"), "--dump-plan", str(plan), "--quiet"])

        slides = json.loads(plan.read_text(encoding="utf-8"))["slides"]
        # 記事は段落が文ごとの箇条書きになる(シナリオでは分けない)。
        assert [b["text"] for b in slides[-1]["bullets"]] == ["一つ目の文です。", "二つ目の文です。"]


class TestFromUrl:
    def test_a_public_article_becomes_a_presentation(self, monkeypatch, tmp_path):
        fake_network(monkeypatch, note=note_payload(body=BODY))
        out = tmp_path / "article.pptx"

        assert cli.main([ARTICLE_URL, "-o", str(out), "--quiet"]) == cli.EXIT_OK

        prs = Presentation(str(out))
        assert prs.core_properties.title == "トークンとは？"
        assert len(prs.slides) >= 3

    def test_the_images_are_saved_next_to_the_presentation(self, monkeypatch, tmp_path):
        fake_network(monkeypatch, note=note_payload(body=BODY))
        out = tmp_path / "article.pptx"

        cli.main([ARTICLE_URL, "-o", str(out), "--quiet"])

        assert (tmp_path / "article_images" / "image_001.png").is_file()

    def test_the_image_directory_can_be_chosen(self, monkeypatch, tmp_path):
        fake_network(monkeypatch, note=note_payload(body=BODY))

        cli.main(
            [
                ARTICLE_URL,
                "-o",
                str(tmp_path / "article.pptx"),
                "--image-dir",
                str(tmp_path / "assets"),
                "--quiet",
            ]
        )

        assert (tmp_path / "assets" / "image_001.png").is_file()

    def test_the_output_name_defaults_to_the_article_key(self, monkeypatch, tmp_path):
        fake_network(monkeypatch, note=note_payload(body=BODY))
        monkeypatch.chdir(tmp_path)

        assert cli.main([ARTICLE_URL, "--quiet"]) == cli.EXIT_OK
        assert (tmp_path / "na01bf3bed64d.pptx").is_file()

    def test_what_was_taken_in_can_be_written_out(self, monkeypatch, tmp_path):
        fake_network(monkeypatch, note=note_payload(body=BODY))
        dump = tmp_path / "article.json"

        cli.main([ARTICLE_URL, "-o", str(tmp_path / "a.pptx"), "--dump-article", str(dump), "--quiet"])

        data = json.loads(dump.read_text(encoding="utf-8"))
        assert data["url"] == ARTICLE_URL
        assert [b["type"] for b in data["blocks"]] == ["paragraph", "heading", "paragraph", "image"]
        assert data["images"][0]["width"] == 1

    def test_the_image_keeps_its_place_in_the_article(self, monkeypatch, tmp_path):
        """図の前後の本文が、記事と同じ並びでスライドになること。"""
        fake_network(monkeypatch, note=note_payload(body=BODY))
        plan = tmp_path / "plan.json"

        cli.main([ARTICLE_URL, "-o", str(tmp_path / "a.pptx"), "--dump-plan", str(plan), "--quiet"])

        kinds = [s["kind"] for s in json.loads(plan.read_text(encoding="utf-8"))["slides"]]
        assert kinds == ["title", "bullets", "bullets", "image"]

    def test_a_failure_reports_the_reason_and_writes_nothing(self, monkeypatch, tmp_path, capsys):
        fake_network(monkeypatch, status=404)
        out = tmp_path / "article.pptx"

        assert cli.main([ARTICLE_URL, "-o", str(out), "--quiet"]) == cli.EXIT_FETCH_FAILED

        assert "HTTP 404" in capsys.readouterr().err
        assert not out.exists()

    def test_a_url_from_another_site_is_refused_before_fetching(self, monkeypatch, tmp_path, capsys):
        requested = fake_network(monkeypatch, note=note_payload(body=BODY))

        assert cli.main(["https://example.com/a", "--quiet"]) == cli.EXIT_USAGE

        assert requested == []
        assert "note.com" in capsys.readouterr().err

    def test_an_existing_output_is_not_overwritten_without_force(self, monkeypatch, tmp_path):
        fake_network(monkeypatch, note=note_payload(body=BODY))
        out = tmp_path / "article.pptx"
        out.write_bytes(b"")

        assert cli.main([ARTICLE_URL, "-o", str(out), "--quiet"]) == cli.EXIT_EXISTS
        assert out.read_bytes() == b""
