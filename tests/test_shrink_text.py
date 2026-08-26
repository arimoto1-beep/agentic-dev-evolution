"""入りきらない文字を小さくする処理が、実際に効いているかを確かめる。

以前は `a:normAutofit` の `fontScale`(「この枠の文字はこの割合で縮めて表示する」
という指示)を書いていた。PowerPoint はこれを守るが、**LibreOffice は読み飛ばす**。
この環境はスライド画像も動画も LibreOffice を通して作るため、資料の上では収まって
いるのに画像でははみ出したまま、という食い違いが起きていた(本文がページ番号の帯に
重なる)。どちらが本当かは、資料と画像の両方を開いて見比べるまで分からない。

そこで縮小率を指示として書くのをやめ、**文字の大きさそのもの** を書き換えている。
ここで確かめたいのは次の 3 つ。

* 縮小が、道具に解釈を任せる指示ではなく、実際の文字の大きさとして書かれること
* 縮めた結果が、本文に使える範囲(= ページ番号の帯の手前)に収まること
* 縮める必要がないときは、何も小さくしないこと
"""

from __future__ import annotations

import os
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt

from note2slides import layout as layout_mod
from note2slides import oxml_utils
from note2slides.renderer import render_deck
from note2slides.scenario import build_deck, parse_scenario
from note2slides.style import Style

HEAD = "---\ntype: scenario\ntitle: サンプル\n---\n\n"


def deck_of(body: str):
    return build_deck(parse_scenario(HEAD + body, source_path="lesson.md"))


def rendered(body: str):
    """シナリオを資料にして、書き出したものを読み直す。

    縮小は書き出しのときに決まるので、資料として保存したものを見る。
    """
    with tempfile.TemporaryDirectory() as directory:
        path = os.path.join(directory, "deck.pptx")
        render_deck(deck_of(body), path, Style())
        return Presentation(path)


def text_frames(pptx_slide):
    return [s.text_frame for s in pptx_slide.shapes if s.has_text_frame]


def frame_with(pptx_slide, needle: str):
    for frame in text_frames(pptx_slide):
        if needle in frame.text:
            return frame
    raise AssertionError(f"{needle!r} を含む枠がありません")


def run_sizes(frame):
    return [
        run.font.size.pt
        for paragraph in frame.paragraphs
        for run in paragraph.runs
        if run.font.size is not None
    ]


def a_text_frame():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    return shape.text_frame


# ---------------------------------------------------------------------------
# shrink_text そのもの
# ---------------------------------------------------------------------------


def test_shrink_text_changes_the_actual_font_size():
    frame = a_text_frame()
    run = frame.paragraphs[0].add_run()
    run.text = "あ"
    run.font.size = Pt(20)

    oxml_utils.shrink_text(frame, 0.5)

    assert run.font.size.pt == 10


def test_shrink_text_does_not_leave_a_scale_for_the_viewer_to_apply():
    # fontScale を書くと、PowerPoint は縮めて LibreOffice は縮めない。
    # 同じ資料が道具によって違う見た目になるため、書いてはいけない。
    frame = a_text_frame()
    run = frame.paragraphs[0].add_run()
    run.text = "あ"
    run.font.size = Pt(20)

    oxml_utils.shrink_text(frame, 0.5)

    body_properties = frame._txBody.bodyPr
    autofit = body_properties.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}normAutofit"
    )
    assert autofit is not None, "あとから文字を足したときの自動調整は残す"
    assert autofit.get("fontScale") is None
    assert autofit.get("lnSpcReduction") is None


def test_shrink_text_also_narrows_the_gap_between_paragraphs():
    frame = a_text_frame()
    first = frame.paragraphs[0]
    first.add_run().text = "あ"
    second = frame.add_paragraph()
    second.space_before = Pt(10)
    second.add_run().text = "い"

    oxml_utils.shrink_text(frame, 0.5)

    assert second.space_before.pt == 5


def test_shrink_text_reduces_line_spacing_only_when_it_is_a_multiple():
    frame = a_text_frame()
    ratio = frame.paragraphs[0]
    ratio.line_spacing = 1.5
    ratio.add_run().text = "あ"
    absolute = frame.add_paragraph()
    absolute.line_spacing = Pt(30)
    absolute.add_run().text = "い"

    oxml_utils.shrink_text(frame, 0.5, line_reduction=0.2)

    assert ratio.line_spacing == 1.2
    assert absolute.line_spacing.pt == 30


# ---------------------------------------------------------------------------
# 資料として見たとき
# ---------------------------------------------------------------------------


def test_long_body_is_written_smaller_than_the_standard_size():
    style = Style()
    standard = style.body_size(0)
    long_body = "\n".join(f"- 入りきらないほど長い本文の{n}行目です。" for n in range(1, 26))

    pptx = rendered("## 見出し\n\n### 画面\n\n" + long_body + "\n")

    frame = frame_with(pptx.slides[0], "1行目")
    assert max(run_sizes(frame)) < standard


def test_slightly_too_much_body_ends_up_fitting():
    # はみ出した本文がページ番号の帯に重なるのが、直したかった症状そのもの。
    # 本文に使える範囲はもともと帯の手前で終わっているので、そこへ収まれば重ならない。
    style = Style()
    body = "\n".join(f"- 少しだけ入りきらない本文の{n}行目です。" for n in range(1, 13))
    source = "## 見出し\n\n### 画面\n\n" + body + "\n"
    box = layout_mod.body_box(style)
    bullets = deck_of(source).slides[0].bullets
    natural = layout_mod.bullets_height(bullets, style, box.width_pt)
    assert natural > box.height_pt, "縮小が要る量になっていること(前提の確認)"

    pptx = rendered(source)

    frame = frame_with(pptx.slides[0], "1行目")
    scale = max(run_sizes(frame)) / style.body_size(0)
    assert natural * scale <= box.height_pt + 1.0
    assert style.body_top + box.height <= style.theme.footer_top


def test_body_is_not_shrunk_past_the_point_of_being_readable():
    # いくらでも小さくすれば必ず収まるが、それは読めない資料になるだけなので
    # 下限で止める(ここまで来たら、画面を分けるかどうかは書いた人が決める)。
    style = Style()
    far_too_much = "\n".join(f"- 到底入りきらない本文の{n}行目です。" for n in range(1, 41))

    pptx = rendered("## 見出し\n\n### 画面\n\n" + far_too_much + "\n")

    frame = frame_with(pptx.slides[0], "1行目")
    assert max(run_sizes(frame)) >= style.body_size(0) * 0.6


def test_a_title_too_wide_for_the_slide_is_written_smaller():
    style = Style()
    title = "とても長い見出しでスライドの幅にはまったく収まらないことが分かっているもの" * 2

    pptx = rendered(f"## {title}\n\n### 画面\n\n本文です。\n")

    frame = frame_with(pptx.slides[0], "とても長い見出し")
    assert max(run_sizes(frame)) < style.slide_title_size


def test_code_too_wide_for_the_slide_is_written_smaller():
    style = Style()
    wide = "x = " + " + ".join(f"variable_number_{n}" for n in range(1, 15))

    pptx = rendered("## 見出し\n\n### 画面\n\n```python\n" + wide + "\n```\n")

    frame = frame_with(pptx.slides[0], "variable_number_1")
    assert max(run_sizes(frame)) < style.code_size


def test_content_that_fits_is_left_at_the_standard_size():
    # 縮める必要がないものまで縮めていないこと(縮小が常時かかっていないこと)。
    style = Style()

    pptx = rendered("## 見出し\n\n### 画面\n\n- 短い項目です\n- もう一つの短い項目です\n")

    frame = frame_with(pptx.slides[0], "短い項目")
    assert max(run_sizes(frame)) == style.body_size(0)
