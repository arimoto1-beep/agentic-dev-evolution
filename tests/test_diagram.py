"""図解(流れ・枠)。

教材の図は、これまで ASCII アートをコードブロックに書くしかなかった。文字を
等幅で並べる書き方は「日本語と罫線が同じ幅で描かれる」ことを前提にしていて、
資料 -> PDF -> 画像と変換する間にどこかで崩れる。しかも崩れたことは、画像を
目で見るまで分からない。

そこで、箱と矢印を **図形として** 置く。ここで確かめたいのは次の 3 つ。

* 書いた項目が、書いた順にそのまま図になること(構成を推測しない)
* 図の大きさの見積り(layout)と、実際に描く大きさ(renderer)が一致すること
* ナレーションが、図解 1 つにつき案内文を 1 度だけ作ること
"""

from __future__ import annotations

import pytest
from pptx import Presentation

from note2slides import guidance, layout, narration
from note2slides.model import (
    DIAGRAM_BOUNDARY,
    DIAGRAM_FLOW,
    DIAGRAM_LANES,
    DIAGRAM_STEPS,
    DIAGRAM_FRAME,
    KIND_DIAGRAM,
    SHAPE_DIAGRAM,
    SHAPE_DIAGRAM_ITEM,
    Bullet,
    Content,
    Run,
    Deck,
    Slide,
    boundary_parts,
    diagram_shape_of,
    lane_parts,
    parse_shape_name,
    step_parts,
)
from note2slides.renderer import render_deck
from note2slides.scenario import ScenarioError, build_deck, parse_scenario
from note2slides.style import Style

HEAD = "---\ntype: scenario\ntitle: サンプル\n---\n\n"


def deck_of(body: str):
    return build_deck(parse_scenario(HEAD + body, source_path="lesson.md"))


def screen(block: str) -> str:
    return "## 見出し\n\n### 画面\n\n" + block + "\n\n### ナレーション\n\n説明です。\n"


FLOW = "```流れ\n受け取る\n処理する\n返す\n```"
FRAME = "```枠\n指示\n会話\n質問\n```"


class TestWhatCountsAsADiagram:
    """どう書いたら図解になるか。"""

    @pytest.mark.parametrize(
        "lang, shape",
        [
            ("流れ", DIAGRAM_FLOW),
            ("flow", DIAGRAM_FLOW),
            ("枠", DIAGRAM_FRAME),
            ("frame", DIAGRAM_FRAME),
            ("FLOW", DIAGRAM_FLOW),
        ],
    )
    def test_japanese_and_english_names_both_work(self, lang, shape):
        assert diagram_shape_of(lang) == shape

    @pytest.mark.parametrize("lang", ["", "text", "python", "bash", "図"])
    def test_other_code_blocks_are_still_code(self, lang):
        """図解の名前でないコードブロックは、これまでどおりコードのまま。

        ここが崩れると、既存の資料のコードが図に化ける。
        """
        assert diagram_shape_of(lang) is None

    def test_a_text_block_stays_code_in_a_scenario(self):
        deck = deck_of(screen("```text\nA\nB\n```"))
        assert deck.slides[0].kind != KIND_DIAGRAM


class TestWhatIsWrittenIsWhatIsDrawn:
    """シナリオが正本であること(構成を推測しない)。"""

    def test_each_line_becomes_one_item_in_order(self):
        deck = deck_of(screen(FLOW))
        slide = deck.slides[0]
        assert slide.kind == KIND_DIAGRAM
        assert slide.diagram_shape == DIAGRAM_FLOW
        assert slide.diagram_items == ["受け取る", "処理する", "返す"]

    def test_blank_lines_and_indentation_do_not_add_items(self):
        deck = deck_of(screen("```流れ\n  受け取る  \n\n   返す\n\n```"))
        assert deck.slides[0].diagram_items == ["受け取る", "返す"]

    def test_a_frame_keeps_its_own_shape(self):
        deck = deck_of(screen(FRAME))
        assert deck.slides[0].diagram_shape == DIAGRAM_FRAME

    def test_an_empty_diagram_stops_with_the_place(self):
        """教材では、図が抜けたまま動画にしても直しようがない(図の指定と同じ扱い)。"""
        with pytest.raises(ScenarioError) as err:
            deck_of(screen("```流れ\n\n```"))
        assert "lesson.md" in str(err.value)
        assert "図の中身が空です" in str(err.value)

    def test_too_many_items_warns_but_does_not_split(self):
        """分けるかどうかを決めるのは書いた人。勝手に次の画面へ送らない。"""
        deck = deck_of(screen("```流れ\n" + "\n".join(f"手順{i}" for i in range(8)) + "\n```"))
        assert len(deck.slides) == 1
        assert len(deck.slides[0].diagram_items) == 8
        assert any("図の項目が 8 個" in w for w in deck.warnings)

    def test_a_diagram_can_sit_beside_text_on_one_screen(self):
        deck = deck_of(screen("説明の文です。\n\n" + FLOW))
        slide = deck.slides[0]
        assert [part.kind for part in slide.parts][-1] == KIND_DIAGRAM

    def test_the_plan_dump_shows_the_diagram(self):
        """--dump-plan で、図の中身を目で確かめられること。"""
        data = deck_of(screen(FLOW)).to_dict()
        assert data["slides"][0]["diagram"] == {
            "shape": DIAGRAM_FLOW,
            "items": ["受け取る", "処理する", "返す"],
        }


class TestArticleInputToo:
    """記事(教材シナリオでない Markdown)でも同じ書き方が使えること。"""

    def test_a_flow_block_in_an_article_becomes_a_diagram(self):
        from note2slides.markdown_reader import parse_article
        from note2slides.planner import plan_deck

        deck = plan_deck(parse_article("# 題\n\n## 見出し\n\n" + FLOW + "\n"))
        diagrams = [s for s in deck.slides if s.kind == KIND_DIAGRAM]
        assert len(diagrams) == 1
        assert diagrams[0].diagram_items == ["受け取る", "処理する", "返す"]
        assert diagrams[0].title == "見出し"


class TestTheEstimateMatchesTheDrawing:
    """見積り(layout)と実際(renderer)がずれると、図が下の中身に重なる。"""

    @pytest.mark.parametrize("shape", [DIAGRAM_FLOW, DIAGRAM_FRAME])
    @pytest.mark.parametrize("count", [1, 3, 6])
    def test_the_drawing_stays_inside_the_place_it_was_given(self, shape, count, tmp_path):
        """layout が配った場所から、描いた図がはみ出さないこと。

        図解は場所が余っていれば大きく描くので、「見積りと同じ大きさ」では
        なくなった。守るべきなのは **配られた場所に収まること** で、
        ここがずれると図が下の中身に重なる(それは画像を見るまで分からない)。
        """
        style = Style()
        content = Content(
            kind=KIND_DIAGRAM,
            diagram_shape=shape,
            diagram_items=[f"項目{i}" for i in range(count)],
        )
        body = layout.body_box(style)
        assert layout.diagram_height(content, style, body.width) > 0

        path = str(tmp_path / "d.pptx")
        render_deck(Deck(slides=[content.as_slide(title="見出し")], title="題"), path, style)
        outer = [s for s in _diagram_shapes(path) if _kind(s) == SHAPE_DIAGRAM]
        assert len(outer) == 1
        left = outer[0].left / 914400
        top = outer[0].top / 914400
        width = outer[0].width / 914400
        height = outer[0].height / 914400
        assert left >= body.left - 0.02
        assert top >= body.top - 0.02
        assert left + width <= body.left + body.width + 0.02
        assert top + height <= body.top + body.height + 0.02

    def test_a_diagram_is_shrunk_rather_than_overflowing(self, tmp_path):
        """場所が足りないときは小さく描く(下の中身に重ねない)。"""
        style = Style()
        items = [f"項目{i}" for i in range(6)]
        content = Content(kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=items)
        path = str(tmp_path / "d.pptx")
        slide = Slide(
            kind="content",
            title="見出し",
            parts=[Content(kind="bullets", bullets=[]), content],
        )
        # 本文の場所を図と文章で分け合うので、図には全部は渡らない。
        render_deck(Deck(slides=[slide], title="題"), path, style)
        outer = [s for s in _diagram_shapes(path) if _kind(s) == SHAPE_DIAGRAM][0]
        bottom = (outer.top + outer.height) / 914400
        assert bottom <= style.body_top + style.body_height + 0.05

    def test_every_item_is_drawn(self, tmp_path):
        style = Style()
        content = Content(
            kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=["A", "B", "C"]
        )
        path = str(tmp_path / "d.pptx")
        render_deck(Deck(slides=[content.as_slide(title="見出し")], title="題"), path, style)
        texts = [
            s.text_frame.text
            for s in _diagram_shapes(path)
            if _kind(s) == SHAPE_DIAGRAM_ITEM and s.has_text_frame and s.text_frame.text
        ]
        assert texts == ["A", "B", "C"]


class TestTheNarrationGuidesToTheDiagram:
    """図解は画面に出ているが、文字としては読み上げられない(表・コードと同じ)。"""

    def test_a_flow_is_read_in_the_order_on_screen(self):
        text = guidance.describe_diagram(DIAGRAM_FLOW, ["文章", "トークン", "回答"])
        assert text == "画面の図をご覧ください。上から順に、文章、トークン、回答と進みます。"

    def test_a_frame_says_what_is_inside(self):
        text = guidance.describe_diagram(DIAGRAM_FRAME, ["指示", "会話"])
        assert text == "画面の図をご覧ください。枠の中には、指示、会話が入っています。"

    def test_nothing_is_added_beyond_what_is_on_screen(self):
        """画面に無い語を足さないこと(表・コードと同じ約束)。"""
        text = guidance.describe_diagram(DIAGRAM_FLOW, ["文章", "トークン"])
        for word in ("文章", "トークン"):
            assert word in text
        assert "つまり" not in text and "重要" not in text

    def test_the_guidance_is_made_once_per_diagram(self, tmp_path):
        """図解は複数の図形でできている。案内文が項目の数だけ出ては困る。"""
        content = Content(
            kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=["A", "B", "C"]
        )
        path = str(tmp_path / "d.pptx")
        render_deck(Deck(slides=[content.as_slide(title="見出し")], title="題"), path, Style())
        lines = narration.extract_script(path).segments
        assert len(lines) == 1
        assert lines[0].text.count("画面の図をご覧ください") == 1
        assert "左から順に、A、B、Cと進みます" in lines[0].text  # 3 項目は横に並ぶ

    def test_the_items_are_not_read_twice(self, tmp_path):
        """案内文で読んだ項目を、本文としてもう一度読まないこと。

        図解は文字を持つ図形の集まりなので、本文を集める側から外しておかないと
        「上から順に、受け取る、返すと進みます。受け取る。返す。」になる。
        """
        content = Content(
            kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=["受け取る", "返す"]
        )
        path = str(tmp_path / "d.pptx")
        render_deck(Deck(slides=[content.as_slide(title="見出し")], title="題"), path, Style())
        text = narration.extract_script(path).segments[0].text
        assert text.count("受け取る") == 1
        assert text.count("返す") == 1

    def test_two_diagrams_on_one_screen_do_not_mix(self, tmp_path):
        flow = Content(kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=["A", "B"])
        frame = Content(kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FRAME, diagram_items=["X", "Y"])
        path = str(tmp_path / "d.pptx")
        render_deck(
            Deck(slides=[Slide(kind="content", title="見出し", parts=[flow, frame])], title="題"),
            path,
            Style(),
        )
        text = narration.extract_script(path).segments[0].text
        assert "左から順に、A、Bと進みます" in text  # 短い流れは横に並ぶ
        assert "枠の中には、X、Yが入っています" in text
        assert "A、B、X、Y" not in text

    def test_a_written_narration_is_used_instead(self, tmp_path):
        """教材シナリオでナレーションを書いた画面では、案内文を組み立てない。"""
        deck = build_deck(parse_scenario(HEAD + screen(FLOW), source_path="lesson.md"))
        path = str(tmp_path / "d.pptx")
        render_deck(deck, path, Style())
        assert narration.extract_script(path).segments[0].text == "説明です。"

    def test_time_is_left_to_look_at_the_diagram(self):
        assert guidance.hold_for_diagram(["A"]) > 0
        assert guidance.hold_for_diagram(["A", "B", "C"]) > guidance.hold_for_diagram(["A"])
        assert guidance.hold_for_diagram([f"項目{i}" for i in range(20)]) <= 3.0


def _diagram_shapes(path: str):
    prs = Presentation(path)
    return list(prs.slides[0].shapes)


def _kind(shape) -> str:
    kind, _, _ = parse_shape_name(getattr(shape, "name", ""))
    return kind


class TestTheDiagramUsesThePlaceItIsGiven:
    """図解は絵なので、渡された場所を使う。

    gen28 までは、大きさが **中身の文字の長さだけ** で決まっていた。短い語の
    流れ図は箱の幅が下限(`diagram_min_width`)に張り付き、画面の左右が大きく
    空いたまま小さく描かれる。44 枚を通して見ると、どの図解も同じ細い列に
    見えるのはこれが理由だった。

    ここで確かめるのは 2 つ。
    * 短い流れは **左から右へ** 並ぶこと(「A から B へ」の自然な読み順)
    * 場所が余っていれば **大きく** 描くこと(ただし見出しより大きくはしない)
    """

    def _outer_and_items(self, path):
        shapes = _diagram_shapes(path)
        outer = [s for s in shapes if _kind(s) == SHAPE_DIAGRAM][0]
        items = [
            s
            for s in shapes
            if _kind(s) == SHAPE_DIAGRAM_ITEM and s.has_text_frame and s.text_frame.text
        ]
        return outer, items

    def _render(self, tmp_path, content, style=None, parts=None):
        style = style or Style()
        path = str(tmp_path / "d.pptx")
        slide = (
            Slide(kind="content", title="見出し", parts=parts)
            if parts
            else content.as_slide(title="見出し")
        )
        render_deck(Deck(slides=[slide], title="題"), path, style)
        return path

    def test_a_short_flow_is_laid_out_left_to_right(self, tmp_path):
        content = Content(
            kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=["AI", "MCP", "AWS"]
        )
        _, items = self._outer_and_items(self._render(tmp_path, content))

        assert [i.text_frame.text for i in items] == ["AI", "MCP", "AWS"]
        # 左から右へ。上下は同じ位置に並ぶ。
        assert items[0].left < items[1].left < items[2].left
        assert items[0].top == items[1].top == items[2].top

    def test_a_long_flow_stays_stacked(self, tmp_path):
        """横に並べるのは、そのままの文字で収まるときだけ。

        収まらない長さの項目を無理に横へ並べると、箱の中で文字が折り返され、
        1 つ 1 つが読めなくなる。それなら縦に積むほうがよい。
        """
        items = [
            "試験データを切り替える",
            "Step Functionsを実行する",
            "結果を確認する",
            "ログを確認する",
            "試験結果を判定する",
        ]
        content = Content(kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=items)
        _, drawn = self._outer_and_items(self._render(tmp_path, content))

        assert [i.text_frame.text for i in drawn] == items
        assert drawn[0].top < drawn[1].top
        assert drawn[0].left == drawn[1].left

    def test_a_frame_is_never_laid_out_left_to_right(self, tmp_path):
        """枠図は「中に何が入っているか」の図で、順番の図ではない。"""
        content = Content(
            kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FRAME, diagram_items=["指示", "会話"]
        )
        _, items = self._outer_and_items(self._render(tmp_path, content))

        assert items[0].top < items[1].top

    def test_the_narration_follows_how_it_was_actually_drawn(self, tmp_path):
        """画面が左から右なら「左から順に」と言うこと。

        画面と言っていることが食い違うと、聞いている側だけが混乱する。
        向きを決めるのは layout なので、ナレーションはその結果を見る。
        """
        across = Content(
            kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=["AI", "MCP", "AWS"]
        )
        down = Content(
            kind=KIND_DIAGRAM,
            diagram_shape=DIAGRAM_FLOW,
            diagram_items=[
                "試験データを切り替える",
                "Step Functionsを実行する",
                "結果を確認する",
                "ログを確認する",
                "試験結果を判定する",
            ],
        )
        tmp_path.joinpath("a").mkdir()
        tmp_path.joinpath("b").mkdir()
        a = narration.extract_script(self._render(tmp_path / "a", across)).segments[0].text
        b = narration.extract_script(self._render(tmp_path / "b", down)).segments[0].text

        assert "左から順に" in a and "上から順に" not in a
        assert "上から順に" in b and "左から順に" not in b

    def test_a_diagram_alone_on_a_screen_is_drawn_larger(self, tmp_path):
        """1 枚を図解だけに使うなら、その場所を使って大きく描く。"""
        style = Style()
        content = Content(
            kind=KIND_DIAGRAM,
            diagram_shape=DIAGRAM_FLOW,
            diagram_items=["機械判定", "AIレビュー", "最終判定"],
        )
        natural = layout.diagram_geometry(content, style, layout.body_box(style).width)
        outer, _ = self._outer_and_items(self._render(tmp_path, content, style))

        assert outer.height / 914400 > natural.height * 1.2

    def test_the_diagram_text_never_gets_bigger_than_the_slide_title(self, tmp_path):
        """図の中の語が、その画面の題より目立つことにならないようにする。"""
        style = Style()
        content = Content(
            kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=["A", "B"]
        )
        body = layout.body_box(style)
        geometry = layout.diagram_geometry(content, style, body.width, body.height)

        assert geometry.font_pt <= style.slide_title_size + 0.01
        assert geometry.font_pt > style.diagram_size  # それでも大きくはなっている

    def test_a_diagram_sharing_the_screen_does_not_take_the_whole_place(self, tmp_path):
        """文章と並ぶ場合は、渡された高さの中に収まる(下の中身を押し出さない)。"""
        style = Style()
        content = Content(
            kind=KIND_DIAGRAM, diagram_shape=DIAGRAM_FLOW, diagram_items=["AI", "MCP", "AWS"]
        )
        text = Content(kind="bullets", bullets=[Bullet(runs=[Run("説明の文です。" * 10)])])
        path = self._render(tmp_path, content, style, parts=[content, text])
        outer, _ = self._outer_and_items(path)

        placed = layout.fit([content, text], style, layout.body_box(style))
        box = placed.parts[0].box
        assert outer.height / 914400 <= box.height + 0.02

    def test_the_leftover_place_goes_to_the_diagram(self, tmp_path):
        """図と文章が並ぶ画面で、余った場所は図に渡る。

        gen28 までは最後の中身(この場合は文章)に渡していた。文章は広い箱を
        もらっても大きくならないので、**余りはどこにも使われず**、図は小さい
        ままだった。
        """
        style = Style()
        diagram = Content(
            kind=KIND_DIAGRAM,
            diagram_shape=DIAGRAM_FLOW,
            diagram_items=["ローカルへ証跡保存", "機械的な期待値チェック", "AIによる解析"],
        )
        text = Content(kind="bullets", bullets=[Bullet(runs=[Run("一行の説明です。")])])
        body = layout.body_box(style)

        placed = layout.fit([diagram, text], style, body)

        natural = layout.diagram_geometry(diagram, style, body.width).height
        assert placed.parts[0].box.height > natural * 1.2

    def test_the_last_text_keeps_room_for_one_more_line(self, tmp_path):
        """全部は渡さない。文章の折り返しが 1 行増えても収まるだけは残す。

        余りを全部図に渡すと、文章がちょうど下端まで来て、ページ番号の帯と
        同じ高さに並ぶ(gen28 が直したのがこの見え方)。
        """
        style = Style()
        diagram = Content(
            kind=KIND_DIAGRAM,
            diagram_shape=DIAGRAM_FLOW,
            diagram_items=["ローカルへ証跡保存", "機械的な期待値チェック", "AIによる解析"],
        )
        text = Content(kind="bullets", bullets=[Bullet(runs=[Run("一行の説明です。")])])
        body = layout.body_box(style)

        placed = layout.fit([diagram, text], style, body)

        estimate = layout.content_height(text, style, body.width)
        one_line = style.line_height_pt(style.body_size(0)) / 72.0
        assert placed.parts[1].box.height >= estimate + one_line - 0.01


BOUNDARY = (
    "```境界\n"
    "人間が、何を作るかを決める\n"
    "↓ 決まった仕様\n"
    "↑ AIだけで決められないこと\n"
    "AIが、どう実現するかを作る\n"
    "```"
)


class TestTheBoundaryDiagram:
    """境界図。1 本の線で上下に分け、線をまたぐものを矢印で示す。

    流れ・枠で書けるのは「順に進む」「中に入っている」の 2 つで、
    **役割の分かれ目** は書けなかった。表に「誰が / 何を受け持つか」と
    書くことはできるが、表には線が無い。「線を引く」と言っている画面に
    線が無いと、言葉と絵が食い違う。

    ここで確かめたいのは 4 つ。

    * `↓` / `↑` の行が線の位置になること(区切りの記号を別に決めない)
    * 線が図の主役として描かれること(箱より広く、箱に隠されない)
    * 矢印が線を **突き抜ける** こと(触れているだけでは、またいで見えない)
    * 案内文が、上下と向きを言い分けること
    """

    @pytest.mark.parametrize("lang", ["境界", "boundary"])
    def test_a_boundary_block_becomes_a_boundary_diagram(self, lang):
        assert diagram_shape_of(lang) == DIAGRAM_BOUNDARY

    def test_the_arrow_lines_decide_where_the_line_is(self):
        parts = boundary_parts(["人間", "↓ 仕様", "↑ 決められないこと", "AI"])
        assert parts.upper == ["人間"]
        assert parts.lower == ["AI"]
        assert [(c.down, c.label) for c in parts.crossings] == [
            (True, "仕様"),
            (False, "決められないこと"),
        ]

    def test_the_crossings_are_not_drawn_as_boxes(self):
        """またぐものは矢印と札になる。箱として数えると場所の見積りがずれる。"""
        content = Content(
            kind=KIND_DIAGRAM,
            diagram_shape=DIAGRAM_BOUNDARY,
            diagram_items=["人間", "↓ 仕様", "AI"],
        )
        assert layout.diagram_items(content) == ["人間", "AI"]

    @pytest.mark.parametrize(
        "block, message",
        [
            ("```境界\n人間\nAI\n```", "線をまたぐもの"),
            ("```境界\n人間\n↓ 仕様\nAI\n↑ 戻す\n```", "離れています"),
            ("```境界\n↓ 仕様\nAI\n```", "線の上に何もありません"),
            ("```境界\n人間\n↓ 仕様\n```", "線の下に何もありません"),
        ],
    )
    def test_a_boundary_that_has_no_line_is_refused(self, block, message):
        """線が決まらない書き方は止める。

        黙って描くと **線の無い境界図** が資料に出る。図の崩れは画像を目で
        見るまで分からない(gen27)ので、書いた時点で知らせる。
        """
        with pytest.raises(ScenarioError) as error:
            deck_of(screen(block))
        assert message in str(error.value)

    def test_the_line_is_wider_than_the_boxes(self, tmp_path):
        """線は箱より左右へはみ出す。

        箱の幅ぴったりで止めると「箱と箱をつなぐ線」に見えて、
        越える・越えないの話に見えない。
        """
        rule, boxes, _ = _boundary_shapes(self._render(tmp_path))
        assert rule.left < min(b.left for b in boxes)
        assert rule.left + rule.width > max(b.left + b.width for b in boxes)

    def test_the_arrows_go_through_the_line(self, tmp_path):
        """矢印は線の上下へ突き抜ける。触れているだけでは、またいで見えない。"""
        rule, _, crossings = _boundary_shapes(self._render(tmp_path))
        assert crossings, "またぐ矢印が描かれていない"
        for arrow in crossings:
            assert arrow.top < rule.top
            assert arrow.top + arrow.height > rule.top + rule.height

    def test_the_label_sits_on_the_side_it_comes_from(self, tmp_path):
        """下りるものの札は線の上、戻るものの札は線の下。

        札を線の上に載せると、そこだけ線が途切れて見え、**1 本の線** という
        図の主題が消える。向きは札の位置で表す。
        """
        path = self._render(tmp_path)
        rule, _, _ = _boundary_shapes(path)
        labels = {language: shape for shape, language in _named_labels(path)}
        assert set(labels) == {"down", "up"}
        assert labels["down"].top + labels["down"].height <= rule.top + rule.height
        assert labels["up"].top >= rule.top

    def test_the_drawing_stays_inside_the_place_it_was_given(self, tmp_path):
        """描いた図形が、外枠(layout が配った場所)からはみ出さないこと。"""
        style = Style()
        body = layout.body_box(style)
        shapes = _diagram_shapes(self._render(tmp_path))
        outer = [s for s in shapes if _kind(s) == SHAPE_DIAGRAM][0]
        assert outer.left / 914400 >= body.left - 0.02
        assert outer.top / 914400 >= body.top - 0.02
        assert (outer.left + outer.width) / 914400 <= body.left + body.width + 0.02
        assert (outer.top + outer.height) / 914400 <= body.top + body.height + 0.02
        slack = 914400 * 0.02
        for shape in shapes:
            if _kind(shape) != SHAPE_DIAGRAM_ITEM:
                continue
            assert shape.left >= outer.left - slack
            assert shape.left + shape.width <= outer.left + outer.width + slack
            assert shape.top >= outer.top - slack
            assert shape.top + shape.height <= outer.top + outer.height + slack

    def test_the_narration_says_which_side_and_which_way(self):
        text = guidance.describe_diagram(
            DIAGRAM_BOUNDARY,
            ["人間が決める", "↓ 決まった仕様", "↑ 決められないこと", "AIが作る"],
        )
        assert text == (
            "画面の図をご覧ください。"
            "線の上は、人間が決める。線の下は、AIが作る。"
            "上から下へ渡るのは、決まった仕様です。"
            "下から上へ戻るのは、決められないことです。"
        )

    def test_the_direction_survives_the_presentation(self, tmp_path):
        """資料に書き出して読み直しても、どちらの向きだったかが残ること。

        案内文を作るのは資料を読む側(narration)なので、向きが資料に
        残っていないと「上から下へ」を言えなくなる。
        """
        path = self._render(tmp_path)
        prs = Presentation(path)
        parts, _ = narration._screen_guidance(prs.slides[0])
        assert parts and "上から下へ渡るのは、決まった仕様です。" in parts[0]
        assert "下から上へ戻るのは、AIだけで決められないことです。" in parts[0]

    def _render(self, tmp_path) -> str:
        deck = deck_of(screen(BOUNDARY))
        path = str(tmp_path / "b.pptx")
        render_deck(deck, path, Style())
        return path


def _named_labels(path: str):
    """境界図の札(向きが名前に残っているもの)を返す。"""
    found = []
    for shape in _diagram_shapes(path):
        kind, _, language = parse_shape_name(getattr(shape, "name", ""))
        if kind == SHAPE_DIAGRAM_ITEM and language in ("down", "up"):
            found.append((shape, language))
    return found


def _boundary_shapes(path: str):
    """境界図の図形を「線・箱・またぐ矢印」に分ける。

    線と矢印は文字を持たない図形で、線は横に長く、矢印は縦に長い。
    """
    labels = {id(s) for s, _ in _named_labels(path)}
    rule = None
    boxes = []
    crossings = []
    for shape in _diagram_shapes(path):
        if _kind(shape) != SHAPE_DIAGRAM_ITEM or id(shape) in labels:
            continue
        text = shape.text_frame.text if shape.has_text_frame else ""
        if text.strip():
            boxes.append(shape)
        elif shape.width > shape.height:
            rule = shape
        else:
            crossings.append(shape)
    return rule, boxes, crossings


LANES = (
    "```レーン\n"
    "人間が決める: 仕様を決めて、承認する\n"
    "AIが作る: 設計する\n"
    "AIが作る: 実装とテストを書く\n"
    "↑ AIだけで決められないこと\n"
    "人間が決める: 受け入れを判断する\n"
    "```"
)


class TestTheLanesDiagram:
    """レーン図。左右 2 列に分け、上から下へ 1 手順ずつ進める。

    流れ・枠・境界のどれでも書けなかったのは **順番と担当を同時に言うこと**
    だった。流れは順番しか言えず、境界は担当しか言えない。全体像の 1 枚は
    「誰が、どの順で、どこで戻すのか」を 1 つの絵で言う必要がある。

    ここで確かめたいのは 5 つ。

    * 書いた順が、そのまま上から下の並びになること
    * 1 行に置く箱が 1 つだけであること(反対の列のその行が空くので、
      またぐ矢印と戻りの横棒が箱の上を通らない)
    * 担当が変わるところで、矢印が縦線を **またぐ** こと
    * 戻りが、**戻り先の箱まで** 届くこと(上を向いたまま終わらない)
    * 案内文が、担当が変わったときだけレーン名を言うこと
    """

    @pytest.mark.parametrize("lang", ["レーン", "lanes", "LANES"])
    def test_a_lanes_block_becomes_a_lanes_diagram(self, lang):
        assert diagram_shape_of(lang) == DIAGRAM_LANES

    def test_the_lane_names_come_from_the_order_they_appear(self):
        """どちらを左に置くかを別に書かせない。先に書いたレーンが左。"""
        parts = lane_parts(["AI: 作る", "人間: 決める", "AI: 直す"])
        assert parts.lanes == ["AI", "人間"]
        assert [(s.lane, s.text) for s in parts.steps] == [
            ("AI", "作る"),
            ("人間", "決める"),
            ("AI", "直す"),
        ]

    @pytest.mark.parametrize("mark", ["：", ":"])
    def test_the_colon_can_be_full_width_or_half_width(self, mark):
        """日本語で書いていて、コロンだけ半角にするのは間違えやすい。"""
        parts = lane_parts([f"人間{mark} 決める", f"AI{mark} 作る"])
        assert parts.lanes == ["人間", "AI"]
        assert parts.steps[0].text == "決める"

    def test_a_colon_inside_the_step_is_kept(self):
        """切るのは最初のコロンだけ。手順の文にコロンがあっても構わない。"""
        parts = lane_parts(["人間: 決める: そして承認する"])
        assert parts.steps[0].text == "決める: そして承認する"

    def test_the_return_remembers_which_step_it_leaves_from(self):
        parts = lane_parts(["人間: 決める", "AI: 作る", "↑ 決められないこと"])
        assert [(r.after, r.label) for r in parts.returns] == [(1, "決められないこと")]

    def test_the_return_is_not_drawn_as_a_box(self):
        """戻りは矢印と札になる。箱として数えると場所の見積りがずれる。"""
        content = Content(
            kind=KIND_DIAGRAM,
            diagram_shape=DIAGRAM_LANES,
            diagram_items=["人間: 決める", "AI: 作る", "↑ 戻す"],
        )
        assert layout.diagram_items(content) == ["決める", "作る"]

    @pytest.mark.parametrize(
        "block, message",
        [
            ("```レーン\n人間: 決める\n作る\n```", "誰の手順か書かれていない"),
            ("```レーン\n人間: 決める\n↓ 仕様\nAI: 作る\n```", "↓ の行があります"),
            ("```レーン\n人間: 決める\n人間: 承認する\n```", "レーンが 1 つ"),
            ("```レーン\n人間: 決める\nAI: 作る\n上司: 見る\n```", "レーンが 3 つ"),
            ("```レーン\n↑ 戻す\n人間: 決める\nAI: 作る\n```", "手順より前"),
        ],
    )
    def test_a_lanes_diagram_that_cannot_be_drawn_is_refused(self, block, message):
        """図が決まらない書き方は、資料にする前に止める。

        レーンが 1 つの図は流れ図と同じもので、担当の分かれ目が無い。それが
        「全体像」として資料に出ても、**画像を目で見るまで気付けない**
        (gen27 以降の前提)。
        """
        with pytest.raises(ScenarioError) as error:
            deck_of(screen(block))
        assert message in str(error.value)

    def test_one_row_holds_one_box(self, tmp_path):
        """1 行に置く箱は 1 つだけ。

        これが崩れると、レーンをまたぐ矢印と戻りの横棒が箱の上を通る。
        図が読めなくなるのは描いてからで、画像を見るまで分からない。
        """
        boxes = _lane_boxes(self._render(tmp_path))
        tops = [shape.top for shape in boxes]
        assert len(set(tops)) == len(tops), "同じ行に 2 つ以上の箱がある"

    def test_the_lanes_are_split_by_one_line(self, tmp_path):
        """列を分ける縦線が 1 本あり、左右の箱がその両側に分かれること。"""
        path = self._render(tmp_path)
        rule = _lane_rule(path)
        assert rule is not None, "列を分ける縦線が引かれていない"
        left, right = _lane_columns(path)
        assert max(b.left + b.width for b in left) <= rule.left
        assert min(b.left for b in right) >= rule.left + rule.width

    def test_the_handover_arrow_crosses_the_line(self, tmp_path):
        """担当が変わるところでは、矢印が縦線をまたぐ。

        またがない矢印は「同じ側で次へ進んだ」だけに見える。担当が移った
        ことは、線を越えることでしか見えない。
        """
        path = self._render(tmp_path)
        rule = _lane_rule(path)
        crossing = [
            s
            for s in _lane_plain(path)
            if s.width > s.height
            and s.left < rule.left
            and s.left + s.width > rule.left + rule.width
        ]
        assert len(crossing) >= 2, "縦線をまたぐ矢印が足りない"

    def test_the_return_reaches_the_box_it_goes_back_to(self, tmp_path):
        """戻りは、戻り先の箱の高さまで届く。

        縦棒の先で終わると「上へ戻る」としか言えず、**どの箱へ戻るのか**
        が示されない。
        """
        path = self._render(tmp_path)
        target = min(_lane_boxes(path), key=lambda s: s.top)
        middle = target.top + target.height / 2
        slack = 914400 * 0.05
        reaching = [
            s
            for s in _lane_plain(path)
            # 横に走る棒で、戻り先の箱と同じ高さにあり、箱の縁まで届いているもの。
            # 縦棒はここに入らない(縦に長い)。またぐ矢印も入らない(行と行の
            # あいだにあって、箱の高さに掛からない)。
            if s.width > s.height
            and s.top <= middle <= s.top + s.height
            and (
                abs(s.left + s.width - target.left) < slack
                or abs(s.left - (target.left + target.width)) < slack
            )
        ]
        assert reaching, "戻りの矢印が、戻り先の箱まで届いていない"

    def test_the_return_label_stays_outside_the_columns(self, tmp_path):
        """戻りの札は列の外の帯に置く。列に重ねると手順の文字が読めない。"""
        path = self._render(tmp_path)
        labels = {language: shape for shape, language in _named_labels(path)}
        assert "up" in labels, "戻りの札が描かれていない"
        left, _ = _lane_columns(path)
        assert labels["up"].left + labels["up"].width <= min(b.left for b in left)

    def test_the_drawing_stays_inside_the_place_it_was_given(self, tmp_path):
        """描いた図形が、layout が配った場所からはみ出さないこと。"""
        path = self._render(tmp_path)
        shapes = _diagram_shapes(path)
        outer = [s for s in shapes if _kind(s) == SHAPE_DIAGRAM][0]
        slack = 914400 * 0.02
        for shape in shapes:
            if _kind(shape) != SHAPE_DIAGRAM_ITEM:
                continue
            assert shape.left >= outer.left - slack
            assert shape.left + shape.width <= outer.left + outer.width + slack
            assert shape.top >= outer.top - slack
            assert shape.top + shape.height <= outer.top + outer.height + slack

    def test_the_narration_names_the_lane_only_when_it_changes(self):
        """毎回レーン名を言うと、どこで担当が移ったのかが聞き取れない。"""
        text = guidance.describe_diagram(
            DIAGRAM_LANES,
            [
                "人間が決める: 仕様を決めて、承認する",
                "AIが作る: 設計する",
                "AIが作る: 実装とテストを書く",
                "↑ AIだけで決められないこと",
                "人間が決める: 受け入れを判断する",
            ],
        )
        assert text == (
            "画面の図をご覧ください。"
            "左は人間が決める、右はAIが作るです。"
            "人間が決める側で、仕様を決めて、承認する。"
            "AIが作る側で、設計する、実装とテストを書く。"
            "人間が決める側で、受け入れを判断する。"
            "人間が決める側へ戻るのは、AIだけで決められないことです。"
        )

    def test_the_lanes_survive_the_presentation(self, tmp_path):
        """資料に書き出して読み直しても、どの手順が誰の列だったかが残ること。

        レーン名は **図形の文字** から取る。図形の名前に入れると
        `parse_shape_name` が小文字にするので、「AI」が「ai」と読まれる。
        """
        prs = Presentation(self._render(tmp_path))
        parts, _ = narration._screen_guidance(prs.slides[0])
        assert parts, "案内文が作られていない"
        assert "左は人間が決める、右はAIが作るです。" in parts[0]
        assert "AIが作る側で、設計する、実装とテストを書く。" in parts[0]
        assert "人間が決める側へ戻るのは、AIだけで決められないことです。" in parts[0]

    def _render(self, tmp_path) -> str:
        deck = deck_of(screen(LANES))
        path = str(tmp_path / "lanes.pptx")
        render_deck(deck, path, Style())
        return path


def _lane_boxes(path: str):
    """レーン図の、手順の箱(レーン名の帯と札は含まない)。"""
    found = []
    for shape in _diagram_shapes(path):
        kind, _, language = parse_shape_name(getattr(shape, "name", ""))
        if kind == SHAPE_DIAGRAM_ITEM and language in ("lane-a", "lane-b"):
            found.append(shape)
    return found


def _lane_columns(path: str):
    """手順の箱を、左の列と右の列に分ける。"""
    left, right = [], []
    for shape in _diagram_shapes(path):
        kind, _, language = parse_shape_name(getattr(shape, "name", ""))
        if kind != SHAPE_DIAGRAM_ITEM:
            continue
        if language == "lane-a":
            left.append(shape)
        elif language == "lane-b":
            right.append(shape)
    return left, right


def _lane_plain(path: str):
    """文字を持たない図形(縦線・矢印・戻りの棒)。"""
    found = []
    for shape in _diagram_shapes(path):
        if _kind(shape) != SHAPE_DIAGRAM_ITEM:
            continue
        text = shape.text_frame.text if shape.has_text_frame else ""
        if not text.strip():
            found.append(shape)
    return found


def _lane_rule(path: str):
    """列を分ける縦線。文字を持たず、いちばん背の高い縦長の図形。"""
    tall = [s for s in _lane_plain(path) if s.height > s.width]
    return max(tall, key=lambda s: s.height) if tall else None


STEPS = (
    "```階段\n"
    "Lv1: それっぽい指摘を出す\n"
    "Lv2: 危ない概念に気づく\n"
    "← Haiku4.5\n"
    "Lv3: 具体的な誤りを特定する\n"
    "Lv4: 横断確認する\n"
    "← Sonnet5 effort max\n"
    "```"
)


class TestTheStepsDiagram:
    """階段図。段を上から下へ 1 段ずつずらして積み、到達点を横に置く。

    流れ・枠・境界・レーンのどれでも書けなかったのは **同じ物差しの上で、
    届いた高さが違うこと** だった。流れは「次に進む」としか言えず(全員が
    最後まで進んでしまう)、境界は 2 段しか作れず、レーンは 1 行 1 箱なので
    段とその到達点を同じ行に置けない。

    ここで確かめたいのは 5 つ。

    * 書いた順が、そのまま上から下の並びになること
    * 段が 1 段ずつ **右へずれる** こと(ずれが「深くなること」そのもの)
    * 段と段が蹴込みでつながって見えること(離れた箱に見えない)
    * 到達点が、**その段の縁まで** 届く矢印になること
    * 到達点の札が、段の上に重ならないこと
    """

    @pytest.mark.parametrize("lang", ["階段", "steps", "STEPS"])
    def test_a_steps_block_becomes_a_steps_diagram(self, lang):
        assert diagram_shape_of(lang) == DIAGRAM_STEPS

    def test_the_levels_keep_the_order_they_were_written(self):
        parts = step_parts(["Lv1: 浅い", "Lv2: 深い", "Lv3: もっと深い"])
        assert [(lv.name, lv.text) for lv in parts.levels] == [
            ("Lv1", "浅い"),
            ("Lv2", "深い"),
            ("Lv3", "もっと深い"),
        ]

    @pytest.mark.parametrize("mark", ["：", ":"])
    def test_the_colon_can_be_full_width_or_half_width(self, mark):
        """レーン図と同じ区切りにしてある(覚えることを増やさない)。"""
        parts = step_parts([f"Lv1{mark} 浅い", f"Lv2{mark} 深い"])
        assert [lv.name for lv in parts.levels] == ["Lv1", "Lv2"]
        assert parts.levels[0].text == "浅い"

    def test_the_reach_remembers_which_level_it_stopped_at(self):
        parts = step_parts(["Lv1: 浅い", "Lv2: 深い", "← Haiku4.5", "Lv3: もっと深い"])
        assert [(r.after, r.label) for r in parts.reaches] == [(1, "Haiku4.5")]

    def test_the_reach_is_not_drawn_as_a_box(self):
        """到達点は矢印と札になる。箱として数えると場所の見積りがずれる。"""
        content = Content(
            kind=KIND_DIAGRAM,
            diagram_shape=DIAGRAM_STEPS,
            diagram_items=["Lv1: 浅い", "Lv2: 深い", "← Haiku4.5"],
        )
        assert layout.diagram_items(content) == ["浅い", "深い"]

    @pytest.mark.parametrize(
        "block, message",
        [
            ("```階段\nLv1: 浅い\n深い\n```", "段の名前が書かれていない"),
            ("```階段\nLv1: 浅い\n↓ 落ちる\nLv2: 深い\n```", "↓ ↑ の行があります"),
            ("```階段\nLv1: 浅い\n```", "段が 1 つしかありません"),
            ("```階段\n← Haiku4.5\nLv1: 浅い\nLv2: 深い\n```", "段より前にあります"),
            ("```階段\nLv1: 浅い\n← A\n← B\nLv2: 深い\n```", "到達点が 2 つ"),
        ],
    )
    def test_a_steps_diagram_that_cannot_be_drawn_is_refused(self, block, message):
        """図が決まらない書き方は、資料にする前に止める。

        段が 1 つの図は階段ではなく、同じ段に到達点が 2 つ付くと札が重なる。
        どちらも黙って描くと **画像を目で見るまで気付けない**
        (gen27 以降の前提)。
        """
        with pytest.raises(ScenarioError) as error:
            deck_of(screen(block))
        assert message in str(error.value)

    def test_each_level_steps_further_right_and_further_down(self, tmp_path):
        """1 段ずつ右へ・下へずれること。

        ずれが 0 なら、ただ縦に積まれた箱で、深くなっていくことが見えない。
        """
        badges = _step_badges(self._render(tmp_path))
        assert len(badges) == 4
        for before, after in zip(badges, badges[1:]):
            assert after.left > before.left, "段が右へずれていない"
            assert after.top > before.top, "段が下へずれていない"

    def test_the_levels_are_joined_by_a_riser(self, tmp_path):
        """段と段のあいだに蹴込みがあること。

        これが無いと、ずらして置いた箱が別々に浮いて見える。階段として
        読めるかどうかは、段そのものではなくここで決まる。
        """
        path = self._render(tmp_path)
        badges = _step_badges(path)
        risers = [s for s in _step_plain(path) if s.height > s.width]
        assert len(risers) == len(badges) - 1, "蹴込みの数が段のあいだの数と合わない"
        for riser, upper, lower in zip(risers, badges, badges[1:]):
            assert upper.top + upper.height <= riser.top + 1
            assert riser.top + riser.height >= lower.top - 1

    def test_the_reach_arrow_touches_the_level_it_points_at(self, tmp_path):
        """到達点の矢印は、その段の縁から出て、見える長さがある。

        帯の中で止まると、どの段のことなのかが決まらない(レーン図の戻りで
        直したのと同じところ)。長さは **どの段の矢印も** 見る。いちばん深い
        段は帯にいちばん近いので、そこだけが矢じりだけになりやすい ——
        浅い段の矢印だけを見ていると、その一本を見落とす。
        """
        path = self._render(tmp_path)
        edges = [s.left + s.width for s in _step_bodies(path)]
        slack = 914400 * 0.03
        arrows = [
            s
            for s in _step_plain(path)
            if s.width > s.height
            and any(abs(s.left - edge) < slack for edge in edges)
        ]
        assert len(arrows) == 2, "到達点の矢印が、段の縁から出ていない"
        for arrow in arrows:
            assert arrow.width > 914400 * 0.3, "矢印が矢じりだけになっている"

    def test_the_reach_label_stays_off_the_levels(self, tmp_path):
        """到達点の札は段の右の帯に置く。段に重なると段の文字が読めない。"""
        path = self._render(tmp_path)
        labels = _step_shapes(path, "step-reach")
        assert len(labels) == 2, "到達点の札が描かれていない"
        deepest = max(s.left + s.width for s in _step_bodies(path))
        for label in labels:
            assert label.left >= deepest, "札が段に重なっている"

    def test_the_drawing_stays_inside_the_place_it_was_given(self, tmp_path):
        """描いた図形が、layout が配った場所からはみ出さないこと。"""
        path = self._render(tmp_path)
        shapes = _diagram_shapes(path)
        outer = [s for s in shapes if _kind(s) == SHAPE_DIAGRAM][0]
        slack = 914400 * 0.02
        for shape in shapes:
            if _kind(shape) != SHAPE_DIAGRAM_ITEM:
                continue
            assert shape.left >= outer.left - slack
            assert shape.left + shape.width <= outer.left + outer.width + slack
            assert shape.top >= outer.top - slack
            assert shape.top + shape.height <= outer.top + outer.height + slack

    def test_the_narration_says_the_level_name_for_the_reach(self):
        """到達点は段の名前で言う。段の文をもう 1 度読むと、同じ文が 2 回出る。"""
        text = guidance.describe_diagram(
            DIAGRAM_STEPS,
            [
                "Lv1: それっぽい指摘を出す",
                "Lv2: 危ない概念に気づく",
                "← Haiku4.5",
                "Lv3: 具体的な誤りを特定する",
            ],
        )
        assert text == (
            "画面の図をご覧ください。"
            "上から下へ、だんだん深くなります。"
            "Lv1、それっぽい指摘を出す。"
            "Lv2、危ない概念に気づく。"
            "Lv3、具体的な誤りを特定する。"
            "Lv2まで届いたのは、Haiku4.5です。"
        )

    def test_the_steps_survive_the_presentation(self, tmp_path):
        """資料に書き出して読み直しても、段の名前と到達点が残ること。"""
        prs = Presentation(self._render(tmp_path))
        parts, _ = narration._screen_guidance(prs.slides[0])
        assert parts, "案内文が作られていない"
        assert "Lv1、それっぽい指摘を出す。" in parts[0]
        assert "Lv2まで届いたのは、Haiku4.5です。" in parts[0]
        assert "Lv4まで届いたのは、Sonnet5 effort maxです。" in parts[0]

    def _render(self, tmp_path) -> str:
        deck = deck_of(screen(STEPS))
        path = str(tmp_path / "steps.pptx")
        render_deck(deck, path, Style())
        return path


def _step_shapes(path: str, language: str):
    found = []
    for shape in _diagram_shapes(path):
        kind, _, found_language = parse_shape_name(getattr(shape, "name", ""))
        if kind == SHAPE_DIAGRAM_ITEM and found_language == language:
            found.append(shape)
    return found


def _step_badges(path: str):
    """段の名前の札(Lv1 など)。"""
    return _step_shapes(path, "step-badge")


def _step_bodies(path: str):
    """段の中身の箱。"""
    return _step_shapes(path, "step-body")


def _step_plain(path: str):
    """文字を持たない図形(蹴込みと、到達点の矢印)。"""
    found = []
    for shape in _diagram_shapes(path):
        if _kind(shape) != SHAPE_DIAGRAM_ITEM:
            continue
        text = shape.text_frame.text if shape.has_text_frame else ""
        if not text.strip():
            found.append(shape)
    return found
