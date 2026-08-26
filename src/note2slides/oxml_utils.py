"""python-pptx が直接扱わない DrawingML 要素を操作するための補助。

OOXML は子要素の順序が決まっているため、スキーマ順に基づいて挿入する。
"""

from __future__ import annotations

from typing import Optional

from lxml import etree
from pptx.oxml.ns import qn
from pptx.util import Pt

# a:pPr の子要素順(ECMA-376 CT_TextParagraphProperties)
_PPR_SEQ = (
    "a:lnSpc",
    "a:spcBef",
    "a:spcAft",
    "a:buClrTx",
    "a:buClr",
    "a:buSzTx",
    "a:buSzPct",
    "a:buSzPts",
    "a:buFontTx",
    "a:buFont",
    "a:buNone",
    "a:buAutoNum",
    "a:buChar",
    "a:tabLst",
    "a:defRPr",
    "a:extLst",
)

# a:rPr / a:defRPr の子要素順(CT_TextCharacterProperties)
_RPR_SEQ = (
    "a:ln",
    "a:noFill",
    "a:solidFill",
    "a:gradFill",
    "a:blipFill",
    "a:pattFill",
    "a:grpFill",
    "a:effectLst",
    "a:effectDag",
    "a:highlight",
    "a:uLnTx",
    "a:uLn",
    "a:uFillTx",
    "a:uFill",
    "a:latin",
    "a:ea",
    "a:cs",
    "a:sym",
    "a:hlinkClick",
    "a:hlinkMouseOver",
    "a:rtl",
    "a:extLst",
)

# a:bodyPr の子要素順(先頭部分のみ)
_BODYPR_SEQ = (
    "a:prstTxWarp",
    "a:noAutofit",
    "a:normAutofit",
    "a:spAutoFit",
    "a:scene3d",
    "a:sp3d",
    "a:flatTx",
    "a:extLst",
)


def _insert_in_order(parent, element, sequence) -> None:
    tag = etree.QName(element).localname
    index = sequence.index(f"a:{tag}")
    successors = sequence[index + 1 :]
    for child in parent:
        child_tag = etree.QName(child).localname
        if f"a:{child_tag}" in successors:
            child.addprevious(element)
            return
    parent.append(element)


def _replace(parent, tag: str, sequence) -> etree._Element:
    for existing in parent.findall(qn(tag)):
        parent.remove(existing)
    element = parent.makeelement(qn(tag), {})
    _insert_in_order(parent, element, sequence)
    return element


def _clear(parent, tags) -> None:
    for tag in tags:
        for existing in parent.findall(qn(tag)):
            parent.remove(existing)


_BULLET_TAGS = ("a:buFont", "a:buNone", "a:buAutoNum", "a:buChar", "a:buClr")


def set_no_bullet(paragraph) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    _clear(pPr, _BULLET_TAGS)
    _replace(pPr, "a:buNone", _PPR_SEQ)


def set_char_bullet(paragraph, char: str, font: str = "Arial", color: Optional[str] = None) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    _clear(pPr, _BULLET_TAGS)
    if color:
        buClr = _replace(pPr, "a:buClr", _PPR_SEQ)
        srgb = buClr.makeelement(qn("a:srgbClr"), {"val": color})
        buClr.append(srgb)
    _replace(pPr, "a:buFont", _PPR_SEQ).set("typeface", font)
    _replace(pPr, "a:buChar", _PPR_SEQ).set("char", char)


def set_auto_number(paragraph, start_at: Optional[int] = None, fmt: str = "arabicPeriod") -> None:
    pPr = paragraph._p.get_or_add_pPr()
    _clear(pPr, _BULLET_TAGS)
    buAutoNum = _replace(pPr, "a:buAutoNum", _PPR_SEQ)
    buAutoNum.set("type", fmt)
    if start_at is not None:
        buAutoNum.set("startAt", str(max(1, start_at)))


def set_indent(paragraph, margin_left_emu: int, hanging_emu: int = 0) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(margin_left_emu)))
    pPr.set("indent", str(int(-hanging_emu)))


def add_line_break(paragraph) -> None:
    """段落の途中に改行(`<a:br/>`)を入れる。

    記事の中の改行(Markdown の行末 2 空白、note の `<br>`)は、段落を分けずに
    行だけを変える。`<a:t>` の中に改行文字を置いても描画側は空白として扱うため、
    行を変えるには要素として入れる必要がある。
    """
    p = paragraph._p
    br = p.makeelement(qn("a:br"), {})
    end_properties = p.find(qn("a:endParaRPr"))
    if end_properties is not None:
        end_properties.addprevious(br)
    else:
        p.append(br)


def set_run_fonts(run, latin: str, east_asian: str, complex_script: Optional[str] = None) -> None:
    """欧文・日本語それぞれのフォントと、文字の出し方を run 自身に書く。

    ひな型のレイアウトが持っている指定は、何も書かなければそのまま効く。
    たとえば `Section Header` のタイトルは `cap="all"` を持っていて、
    **書いた文字と違う文字が出る**(「試験Runner」->「試験RUNNER」)。
    資料の中身は正しいので、画像を目で見るまで気付かない。
    ここで `cap="none"` を書いて、書いたとおりの文字が出るようにする。
    """
    rPr = run._r.get_or_add_rPr()
    rPr.set("cap", "none")
    _replace(rPr, "a:latin", _RPR_SEQ).set("typeface", latin)
    _replace(rPr, "a:ea", _RPR_SEQ).set("typeface", east_asian)
    _replace(rPr, "a:cs", _RPR_SEQ).set("typeface", complex_script or latin)


def shrink_text(text_frame, font_scale: float, line_reduction: float = 0.0) -> None:
    """はみ出す文字を、実際に小さくする。

    `a:normAutofit` の `fontScale` は「この枠の文字はこの割合で縮めて表示する」
    という指示で、PowerPoint はこれを守る。**LibreOffice は読み飛ばす。**
    そのため縮小率を指示として書くと、資料の上では収まっているのに、そこから
    作るスライド画像と動画でははみ出したまま、という食い違いが起きる
    (ページ番号の帯に本文が重なる)。しかも資料と画像のどちらが本当かは、
    両方を開いて見比べるまで分からない。

    ここでは指示を書かず、**文字の大きさそのもの** を書き換える。どの道具で
    開いても同じ結果になり、収まったかどうかを画像で確かめられる。

    `a:normAutofit` 自体は残す(縮小率は付けない)。あとから PowerPoint で
    文字を足したときは、そこから先を PowerPoint の自動調整に任せる。
    """
    for paragraph in text_frame.paragraphs:
        spacing = paragraph.line_spacing
        # 行間は「文字の大きさの何倍」で指定している場合だけ詰める
        # (pt で指定されている場合は、文字を縮めても詰めない)。
        if line_reduction and isinstance(spacing, float):
            paragraph.line_spacing = spacing * (1.0 - line_reduction)
        if paragraph.space_before is not None:
            paragraph.space_before = Pt(paragraph.space_before.pt * font_scale)
        for run in paragraph.runs:
            if run.font.size is not None:
                run.font.size = Pt(run.font.size.pt * font_scale)
    bodyPr = text_frame._txBody.bodyPr
    _clear(bodyPr, ("a:noAutofit", "a:normAutofit", "a:spAutoFit"))
    _replace(bodyPr, "a:normAutofit", _BODYPR_SEQ)
