"""プレゼンテーション資料からナレーション原稿を取り出す。

原稿は「スライド 1 枚 = セリフ 1 本」で、番号はスライド番号(1 起点)と一致する。
音声ファイルの連番もこの番号に合わせるため、スライドと音声の対応は番号だけで
追える。

文章は資料に書かれている文字をそのまま使う。要約・言い換え・補足の生成は
行わない(資料自体が記事本文をそのまま並べたものなので、原稿も記事の文言の
ままになる)。読み上げ元は次の優先順で選び、どれを使ったかを `source` に残す。

    notes   発表者ノート(note2slides が入れた元の本文)
    screen  画面に出ている表・コード・図の案内(ノートが無い場合、`guidance`)
    body    スライド上の本文(表・コード・図が無い場合)
    title   スライドのタイトルだけ(本文が無い場合)
    none    読み上げる文字が無い(無音として扱う)

表・コード・図は画面には出ているが、そのままでは読み上げる文字にならない。
何も言わずに映すと、視聴者は何を見ればよいのか分からないまま次のスライドへ
進んでしまうため、画面に出ている文字だけを使って案内文を組み立てる
(組み立て方は `guidance.py`)。ここでも記事に無い事実は足さない。

読み上げる内容が資料を作る前から決まっている場合(教材シナリオ →
`scenario.py`)は、ノートに指示行(`[note2slides] ...`)が入る。この行がある
ノートは書いた人の指定として扱い、文章が無ければ画面の案内文を組み立てずに
無音のままにする。

原稿は JSON として書き出し・読み込みできる。読みを直したい場合は書き出した
JSON を編集して、それを入力にして合成する。
"""

from __future__ import annotations

import json
import os
from dataclasses import dataclass, field
from typing import List, Optional, Tuple

from . import guidance
from .model import (
    CROSSING_DOWN,
    CROSSING_UP,
    REACH_MARK,
    SHAPE_CODE,
    SHAPE_DIAGRAM,
    SHAPE_DIAGRAM_ITEM,
    SHAPE_FOOTER,
    parse_shape_name,
)

#: 境界図で線をまたぐものの札に、描画側が付ける名前(renderer._draw_boundary)。
#: レーン図の戻りの札も `up` を使う(どちらも「戻るもの」で、案内文の側では
#: 同じ印から読み取る)。
_CROSSING_BY_NAME = {"down": CROSSING_DOWN, "up": CROSSING_UP}

#: レーン図の図形に、描画側が付ける名前(renderer.PptxRenderer.LANE_*)。
#: レーン名は **図形の文字** から取る。名前に入れると `parse_shape_name` が
#: 小文字にするので、「AI」が「ai」になって読み上げられてしまう。
_LANE_HEAD = "lane-head"
_LANE_COLUMNS = ("lane-a", "lane-b")

#: 階段図の図形に、描画側が付ける名前(renderer.PptxRenderer.STEP_*)。
#: 段の名前もレーン名と同じ理由で **図形の文字** から取る。
_STEP_BADGE = "step-badge"
_STEP_BODY = "step-body"
_STEP_REACH = "step-reach"

#: 読み上げ元の種別。
SOURCE_NOTES = "notes"
SOURCE_SCREEN = "screen"
SOURCE_BODY = "body"
SOURCE_TITLE = "title"
SOURCE_NONE = "none"

#: 1 枚に収まらないスライドのタイトルに planner が付ける目印。画面を見れば
#: 続きだと分かるので、読み上げからは外す(「かっこ つづき」と読まれてしまう)。
CONTINUATION_SUFFIX = "（続き）"

SCRIPT_SUFFIXES = (".json",)
PRESENTATION_SUFFIXES = (".pptx",)

#: 発表者ノートに残す指示行の目印。
#:
#: ノートは「そのスライドで読み上げる文章」なので、文章では表せない指定
#: (読み上げないこと・読み終わったあとに画面を見せる時間)だけをこの行に書く。
#: 教材シナリオ(`scenario.py`)が書き込み、ここで読み取る。資料を PowerPoint で
#: 開いて直接書き足すこともできる。
#:
#:     [note2slides] hold=3
#:     [note2slides] narration=none hold=4
#:
#: この行があるノートは「読み上げる内容が決まっている」ものとして扱い、
#: 文章が空なら、画面の案内文を組み立てずに無音のままにする。
DIRECTIVE_PREFIX = "[note2slides]"
#: 読み上げないことを明示する値。
NARRATION_NONE = "none"


class NarrationError(RuntimeError):
    """原稿の取り出し・読み込みに失敗した場合。"""


@dataclass
class NarrationSegment:
    """1 枚のスライドに対応するセリフ。index はスライド番号(1 起点)。

    `hold` は読み終わったあとに置く無音(秒)。コードや図のように、聞くだけでは
    足りず画面を読む必要がある内容のために取る時間で、そのぶんスライドが長く映る。
    """

    index: int
    text: str = ""
    title: str = ""
    source: str = SOURCE_NONE
    hold: float = 0.0

    @property
    def is_empty(self) -> bool:
        return not self.text.strip()

    def to_dict(self) -> dict:
        data = {
            "index": self.index,
            "title": self.title,
            "source": self.source,
            "text": self.text,
        }
        if self.hold:
            data["hold"] = round(self.hold, 3)
        return data

    @classmethod
    def from_dict(cls, data: dict, fallback_index: int) -> "NarrationSegment":
        if not isinstance(data, dict):
            raise NarrationError(f"セリフの形式が正しくありません: {data!r}")
        index = data.get("index", fallback_index)
        if not isinstance(index, int) or index < 1:
            raise NarrationError(f"index は 1 以上の整数にしてください: {index!r}")
        text = data.get("text", "")
        if not isinstance(text, str):
            raise NarrationError(f"{index} 番のセリフの text が文字列ではありません")
        return cls(
            index=index,
            text=text,
            title=str(data.get("title", "")),
            # 手で書いた原稿にも対応するため、source は無ければ推定する。
            source=str(data.get("source") or (SOURCE_NOTES if text.strip() else SOURCE_NONE)),
            hold=_hold_of(data.get("hold", 0.0), index),
        )


def _hold_of(value: object, index: int) -> float:
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        raise NarrationError(f"{index} 番のセリフの hold が数値ではありません: {value!r}")
    if value < 0:
        raise NarrationError(f"{index} 番のセリフの hold は 0 以上にしてください: {value}")
    return float(value)


@dataclass
class NarrationScript:
    segments: List[NarrationSegment] = field(default_factory=list)
    source: str = ""  # 取り出し元の資料
    warnings: List[str] = field(default_factory=list)

    @property
    def count(self) -> int:
        return len(self.segments)

    def to_dict(self) -> dict:
        return {
            "source": os.path.basename(self.source) if self.source else "",
            "count": self.count,
            "segments": [s.to_dict() for s in self.segments],
            "warnings": self.warnings,
        }

    def write(self, path: str) -> str:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(self.to_dict(), f, ensure_ascii=False, indent=2)
            f.write("\n")
        return os.path.abspath(path)


def read_script(path: str) -> NarrationScript:
    """書き出した原稿(JSON)を読み込む。手で直したものも入力にできる。"""
    try:
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
    except OSError as exc:
        raise NarrationError(f"原稿を開けませんでした: {path}\n{exc}") from exc
    except json.JSONDecodeError as exc:
        raise NarrationError(
            f"原稿が JSON として読めませんでした: {path}\n{exc.lineno} 行 {exc.colno} 文字目: {exc.msg}"
        ) from exc

    if isinstance(data, list):  # segments だけを書いた形も許す
        data = {"segments": data}
    if not isinstance(data, dict):
        raise NarrationError(f"原稿の形式が正しくありません: {path}")

    raw = data.get("segments")
    if not isinstance(raw, list) or not raw:
        raise NarrationError(f"原稿にセリフ(segments)がありません: {path}")

    segments = [NarrationSegment.from_dict(item, i) for i, item in enumerate(raw, start=1)]
    script = NarrationScript(
        segments=segments,
        source=str(data.get("source") or path),
        warnings=[str(w) for w in data.get("warnings", []) if isinstance(w, str)],
    )
    _check_indexes(script, path)
    return script


def _check_indexes(script: NarrationScript, path: str) -> None:
    """番号の重複と欠落を弾く。スライドとの対応が崩れたまま合成しないため。"""
    indexes = [s.index for s in script.segments]
    if sorted(indexes) != list(range(1, len(indexes) + 1)):
        raise NarrationError(
            f"原稿の index が 1 からの連番になっていません: {path}\n"
            f"  見つかった index: {sorted(indexes)}\n"
            "スライド番号と音声の番号を合わせるため、1 起点の通し番号にしてください。"
        )
    script.segments.sort(key=lambda s: s.index)


def extract_script(source: str) -> NarrationScript:
    """資料(.pptx)からナレーション原稿を取り出す。

    非表示スライドは画像化でも出力されないため、原稿からも外して警告する。
    """
    suffix = os.path.splitext(source)[1].lower()
    if suffix in SCRIPT_SUFFIXES:
        return read_script(source)
    if suffix not in PRESENTATION_SUFFIXES:
        known = " / ".join(PRESENTATION_SUFFIXES + SCRIPT_SUFFIXES)
        raise NarrationError(f"未対応の入力形式です: {suffix or source}(扱えるのは {known})")
    if not os.path.isfile(source):
        raise NarrationError(f"入力ファイルが見つかりません: {source}")

    from pptx import Presentation
    from pptx.exc import PythonPptxError

    try:
        presentation = Presentation(source)
        slides = list(presentation.slides)
    except (PythonPptxError, ValueError, KeyError, OSError) as exc:
        raise NarrationError(
            f"資料として読み取れませんでした: {source}\n{type(exc).__name__}: {exc}"
        ) from exc

    script = NarrationScript(source=os.path.abspath(source))
    hidden = 0
    number = 0
    for slide in slides:
        if slide.element.get("show") == "0":
            hidden += 1
            continue
        number += 1
        try:
            script.segments.append(_segment_of(slide, number))
        except (PythonPptxError, ValueError, KeyError, TypeError, AttributeError) as exc:
            # どの 1 枚で止まったのかが分からないと、資料のどこを直せばよいか
            # 調べようがない。番号と見出しを添えて伝える。
            raise NarrationError(
                f"スライド {number} からセリフを取り出せませんでした: {source}\n"
                f"  見出し: {_clean(_title_text(slide)) or '(なし)'}\n"
                f"  {type(exc).__name__}: {exc}"
            ) from exc

    if not script.segments:
        raise NarrationError(f"スライドが 1 枚もありませんでした: {source}")
    if hidden:
        script.warnings.append(
            f"非表示のスライド {hidden} 枚は原稿から外しました(スライド画像にも出力されません)。"
        )
    silent = [s.index for s in script.segments if s.is_empty]
    if silent:
        script.warnings.append(
            "読み上げる文字が無いスライドがあります(無音になります): "
            + ", ".join(str(i) for i in silent)
        )
    return script


def _segment_of(slide, number: int) -> NarrationSegment:
    """1 枚分のセリフを決める。ノートが無い場合は画面に出ているものを説明する。"""
    title = _clean(_title_text(slide))
    spoken_title = _spoken_title(title)
    body, directives = split_directives(_notes_text(slide), index=number)
    notes = _clean(body)
    if directives is not None:
        # 読み上げる内容が決まっているノート。文章が無い場合は、画面の案内文を
        # 組み立てずに無音のままにする(書いた人が「ここは読まない」と決めている)。
        return NarrationSegment(
            number,
            text=notes,
            title=title,
            source=SOURCE_NOTES if notes else SOURCE_NONE,
            hold=directives.hold,
        )
    if notes:
        return NarrationSegment(number, text=notes, title=title, source=SOURCE_NOTES)

    # 表紙や章扉はタイトルだけが本文なので、タイトルも読み上げの対象にする。
    screen, hold = _screen_guidance(slide)
    body = _body_texts(slide)
    spoken = [part for part in [spoken_title] + screen + body if part.strip()]
    text = _clean("\n".join(spoken))

    if screen:
        return NarrationSegment(number, text=text, title=title, source=SOURCE_SCREEN, hold=hold)
    if body:
        return NarrationSegment(number, text=text, title=title, source=SOURCE_BODY)
    if text:
        return NarrationSegment(number, text=text, title=title, source=SOURCE_TITLE)
    # 読むものがタイトルしか無い「続き」のスライドは、目印を外して読み上げる。
    if title:
        return NarrationSegment(
            number, text=_base_title(title), title=title, source=SOURCE_TITLE
        )
    return NarrationSegment(number, text="", title=title, source=SOURCE_NONE)


# ---------------------------------------------------------------------------
# 発表者ノートの指示行
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class Directives:
    """ノートに書かれた、文章では表せない指定。"""

    hold: float = 0.0
    narration: str = ""


def compose_notes(text: str, hold: float = 0.0) -> str:
    """ナレーションを発表者ノートの文字列にする(指示が要る場合だけ行を足す)。"""
    body = (text or "").strip()
    parts = []
    if not body:
        parts.append(f"narration={NARRATION_NONE}")
    if hold:
        parts.append(f"hold={hold:g}")
    if not parts:
        return body
    line = " ".join([DIRECTIVE_PREFIX] + parts)
    return f"{body}\n{line}" if body else line


def split_directives(text: str, index: int = 0) -> Tuple[str, Optional[Directives]]:
    """ノートを「読み上げる文章」と「指示」に分ける。

    指示行が 1 行も無い場合は `None` を返す(記事から作った資料や、人が書いた
    ノートはこちらになる)。
    """
    normalized = (text or "").replace("\v", "\n").replace("\r\n", "\n").replace("\r", "\n")
    body: List[str] = []
    found = False
    hold = 0.0
    narration = ""
    for raw in normalized.split("\n"):
        line = raw.strip()
        if not line.lower().startswith(DIRECTIVE_PREFIX):
            body.append(raw)
            continue
        found = True
        for token in line[len(DIRECTIVE_PREFIX) :].split():
            key, _, value = token.partition("=")
            key = key.strip().lower()
            if key == "hold":
                hold = _directive_seconds(value, index, line)
            elif key == "narration":
                narration = _directive_narration(value, index, line)
            else:
                raise NarrationError(
                    f"スライド {index} のノートに知らない指示があります: {token}\n"
                    f"  {line}\n"
                    "  書けるのは narration=none と hold=<秒> です。"
                )
    return "\n".join(body), (Directives(hold=hold, narration=narration) if found else None)


def _directive_narration(value: str, index: int, line: str) -> str:
    narration = value.strip().lower()
    if narration != NARRATION_NONE:
        raise NarrationError(
            f"スライド {index} のノートの narration に書けるのは"
            f" {NARRATION_NONE} だけです: {value!r}\n  {line}\n"
            "  読み上げる文章は、この行ではなくノートの本文に書いてください。"
        )
    return narration


def _directive_seconds(value: str, index: int, line: str) -> float:
    try:
        seconds = float(value)
    except ValueError:
        raise NarrationError(
            f"スライド {index} のノートの hold が秒数ではありません: {value!r}\n  {line}"
        ) from None
    if seconds < 0:
        raise NarrationError(f"スライド {index} のノートの hold は 0 以上にしてください: {seconds}")
    return seconds


def _screen_guidance(slide) -> Tuple[List[str], float]:
    """画面に出ている表・コード・図を、案内文と画面を読む時間にする。

    「続き」かどうかは図形の名前に残っている(`model.shape_name`)。1 枚に
    収まらず分けた表・コードだけが続きで、見出しが同じというだけの別の表は
    続きにはしない。
    """
    parts: List[str] = []
    hold = 0.0
    diagrams = _diagram_items(slide)
    for position, shape in enumerate(slide.shapes):
        kind, continued, language = parse_shape_name(getattr(shape, "name", ""))
        if kind == SHAPE_DIAGRAM:
            items = diagrams.get(position, [])
            text = guidance.describe_diagram(language, items)
            hold = max(hold, guidance.hold_for_diagram(items))
        elif kind == SHAPE_DIAGRAM_ITEM:
            continue  # 図解の中の図形。案内文は外枠のところで 1 度だけ作る
        elif getattr(shape, "has_table", False):
            header, rows = _table_cells(shape.table)
            text = guidance.describe_table(header, rows, continued=continued)
        elif _is_code(shape):
            code = shape.text_frame.text
            text = guidance.describe_code(code, language, continued=continued)
            hold = max(hold, guidance.hold_for_code(code))
        elif _is_picture(shape):
            text = guidance.describe_image()
            hold = max(hold, guidance.hold_for_image())
        else:
            continue
        if text:
            parts.append(text)
    return parts, hold


def _diagram_items(slide) -> dict:
    """図解の外枠ごとに、その中の項目の文字を集める。

    図解は「外枠 -> 項目 -> 項目 …」の順に置かれる(renderer._draw_diagram)。
    1 枚に図解が 2 つあっても混ざらないよう、外枠が現れるたびに集め直す。

    どの外枠のものかは **並び順** で覚える。python-pptx の図形は、
    `slide.shapes` を見るたびに作り直される入れ物なので、同じ図形でも
    `id()` は同じにならず、たまたま一致することもある(別の図の項目が
    混ざる)。並び順なら、同じ資料を読むかぎり同じものを指す。
    """
    found: dict = {}
    current: Optional[List[str]] = None
    lanes: List[str] = []
    badge = ""
    for position, shape in enumerate(slide.shapes):
        kind, _, language = parse_shape_name(getattr(shape, "name", ""))
        if kind == SHAPE_DIAGRAM:
            current = []
            lanes = []
            badge = ""
            found[position] = current
        elif kind == SHAPE_DIAGRAM_ITEM and current is not None:
            text = getattr(shape, "text_frame", None)
            text = text.text.strip() if text is not None else ""
            if not text:
                continue
            if language == _STEP_BADGE:
                # 階段図の段の名前。段の中身は次に来るので、覚えて結び直す。
                badge = text
            elif language == _STEP_BODY:
                current.append(f"{badge}: {text}" if badge else text)
                badge = ""
            elif language == _STEP_REACH:
                # 到達点は、シナリオに書いたのと同じ印を戻して渡す。
                current.append(f"{REACH_MARK} {text}")
            elif language == _LANE_HEAD:
                # レーン名の帯は項目ではない。手順にレーン名を結び直すために
                # 覚えておく(帯は手順より先に描かれる)。
                lanes.append(text)
            elif language in _LANE_COLUMNS:
                lane = _lane_name(lanes, _LANE_COLUMNS.index(language))
                current.append(f"{lane}: {text}" if lane else text)
            else:
                # 境界図で線をまたぐものは、向きが図形の名前に残っている。
                # 案内文が「上から下へ」「下から上へ」を言い分けられるよう、
                # シナリオに書いたのと同じ印を戻して渡す(model.CROSSING_MARKS)。
                mark = _CROSSING_BY_NAME.get(language, "")
                current.append(f"{mark}{text}" if mark else text)
        elif kind not in (SHAPE_DIAGRAM, SHAPE_DIAGRAM_ITEM):
            current = None
    return found


def _lane_name(lanes: List[str], column: int) -> str:
    """レーン図で、その列のレーン名。帯が読めていなければ空にする。"""
    return lanes[column] if 0 <= column < len(lanes) else ""


def _table_cells(table) -> Tuple[List[str], List[List[str]]]:
    """表を「見出しの行」と「中身の行」に分ける。

    見出しかどうかは資料の設定(first_row)に従う。note2slides が作った資料では、
    Markdown の表の見出し行があるときだけ立っている。
    """
    rows = [[cell.text for cell in row.cells] for row in table.rows]
    if not rows:
        return [], []
    if getattr(table, "first_row", False):
        return rows[0], rows[1:]
    return [], rows


def _is_code(shape) -> bool:
    kind, _, _ = parse_shape_name(getattr(shape, "name", ""))
    return kind == SHAPE_CODE and getattr(shape, "has_text_frame", False)


def _is_diagram_part(shape) -> bool:
    """図解を作っている図形かどうか(外枠・項目・矢印)。"""
    kind, _, _ = parse_shape_name(getattr(shape, "name", ""))
    return kind in (SHAPE_DIAGRAM, SHAPE_DIAGRAM_ITEM)


def _is_decoration(shape) -> bool:
    """資料名・ページ番号のような、読み上げない飾りかどうか。"""
    kind, _, _ = parse_shape_name(getattr(shape, "name", ""))
    return kind == SHAPE_FOOTER


def _is_picture(shape) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except (AttributeError, ValueError, KeyError):
        # 種類を判別できない図形は、画像として扱わない(誤って案内文を足さない)。
        return False


def _is_continuation(title: str) -> bool:
    return title.endswith(CONTINUATION_SUFFIX)


def _spoken_title(title: str) -> str:
    """読み上げ用のタイトル。

    「（続き）」のスライドは、同じ見出しを前のスライドで読み上げているので
    繰り返さない(目印をそのまま読むと「かっこ つづき」になってしまう)。
    """
    return "" if _is_continuation(title) else title


def _base_title(title: str) -> str:
    """画面上の目印を外したタイトル。"""
    if _is_continuation(title):
        return title[: -len(CONTINUATION_SUFFIX)].strip()
    return title


def _title_text(slide) -> str:
    try:
        title = slide.shapes.title
    except (AttributeError, ValueError):
        return ""
    return title.text if title is not None else ""


def _notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    return frame.text if frame is not None else ""


def _body_texts(slide) -> List[str]:
    """タイトル以外の本文テキストを、スライド上の並び順で集める。

    表・画像は文字として読み上げられないため対象外(`_screen_guidance` が
    案内文にする)。コード用の図形も、そのまま読み上げても聞き取れないので外す。
    資料名・ページ番号(フッタ)は見た目のための飾りなので、これも外す。
    図の説明(キャプション)はここで拾い、図の案内文のあとに読み上げる。

    図解(箱と矢印)の中の文字も外す。ここで拾うと、`_screen_guidance` が
    作った案内文のあとに、同じ語がもう一度並んで読み上げられる。
    """
    title = slide.shapes.title
    texts = []
    for shape in slide.shapes:
        if title is not None and shape.element is title.element:
            continue
        if _is_code(shape) or _is_decoration(shape) or _is_diagram_part(shape):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text
        if text.strip():
            texts.append(text)
    return texts


def _clean(text: str) -> str:
    """読み上げ用に整える。文字の削除・追加はせず、空白と改行だけを整理する。

    * 縦タブ(pptx の行内改行)を改行にする
    * 行頭・行末の空白を落とし、空行を詰める
    """
    if not text:
        return ""
    normalized = text.replace("\v", "\n").replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in normalized.split("\n")]
    return "\n".join(line for line in lines if line)
