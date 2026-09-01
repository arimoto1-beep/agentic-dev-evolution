"""スライド構成(Deck)を .pptx として書き出す。

出力は Office Open XML なので、PowerPoint だけでなく LibreOffice Impress でも
そのまま開いて編集できる。プレースホルダを使いつつ 16:9 に合わせて配置し直す。
"""

from __future__ import annotations

from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from . import layout as layout_mod
from . import metrics, oxml_utils, text_wrap
from .layout import Box
from .model import (
    DIAGRAM_BOUNDARY,
    DIAGRAM_LANES,
    DIAGRAM_STEPS,
    DIAGRAM_FLOW_ACROSS,
    DIAGRAM_FRAME,
    KIND_CODE,
    KIND_DIAGRAM,
    KIND_IMAGE,
    KIND_SECTION,
    KIND_TABLE,
    KIND_THUMBNAIL,
    KIND_TITLE,
    NUMBER,
    PLAIN,
    QUOTE,
    SHAPE_CODE,
    SHAPE_DIAGRAM,
    SHAPE_DIAGRAM_ITEM,
    SHAPE_FOOTER,
    SHAPE_TABLE,
    Bullet,
    Content,
    Deck,
    Run,
    Slide,
    boundary_parts,
    lane_parts,
    shape_name,
    step_parts,
)
from .style import SLIDE_HEIGHT_EMU, SLIDE_WIDTH_EMU, Style, inches_to_pt

_LAYOUT_TITLE = 0
_LAYOUT_SECTION = 2
_LAYOUT_TITLE_ONLY = 5
_LAYOUT_BLANK = 6

_BULLET_CHARS = {0: "●", 1: "–", 2: "・", 3: "・"}

# 表紙・章扉の配置(inch)。タイトルは行数に応じて縦に伸びるので、
# 「下端(表紙)」「中心(章扉)」を固定して、そこから箱の大きさを決める。
_COVER_LEFT = 1.0
_COVER_WIDTH = 11.33
_COVER_TITLE_MIN_HEIGHT = 1.6
_COVER_RULE_WIDTH = 2.0
#: 題の下端から、罫線・副題までの間(inch)。
_COVER_RULE_GAP = 0.25
_COVER_SUBTITLE_GAP = 0.5
_SECTION_CENTER = 3.5
_SECTION_RULE_TOP = 4.45
#: 表紙の下端に引く帯(inch)。
_COVER_BAND_HEIGHT = 0.11
#: タイトルを何行まで折り返してよいか。
_COVER_TITLE_MAX_LINES = 3
#: 見出しの下に引く、うすい罫線の太さ(inch)。
_HAIRLINE = 0.018
#: フッタの資料名に使ってよい横幅の割合。
_FOOTER_TITLE_SHARE = 0.6

# サムネイル(1 枚絵)の配置(inch)と、題に使う文字の大きさ(pt)。
# 小さく表示されても読めるよう、収まる範囲でいちばん大きい文字を選ぶ。
_THUMB_LEFT = 1.0
_THUMB_WIDTH = 11.33
_THUMB_TITLE_SIZES = (68, 60, 54, 48, 42, 36)
#: 自然な切れ目では収まらない題を小さくしていくときの刻み(pt)。
_THUMB_TITLE_STEP = 2.0
#: これ以下になったら、一覧では読めない(`thumbnail_title_is_cramped` が知らせる)。
_THUMB_TITLE_MIN_SIZE = 20.0
#: それでも収まらない題のための、崩さないための下限(pt)。読めるかどうかは
#: 上の警告に任せ、**版面からはみ出さないこと** だけをここで守る。
_THUMB_TITLE_HARD_MIN = 6.0
_THUMB_TITLE_MAX_LINES = 3
_THUMB_TITLE_MAX_HEIGHT = 3.9
_THUMB_RULE_WIDTH = 2.2
_THUMB_RULE_HEIGHT = 0.08
_THUMB_LABEL_SIZE = 22.0
_THUMB_LABEL_HEIGHT = 0.58
_THUMB_SUBTITLE_SIZE = 26.0
_THUMB_BAND_HEIGHT = 0.16
_THUMB_RULE_GAP = 0.34
_THUMB_SUBTITLE_GAP = 0.26
_THUMB_LABEL_GAP = 0.42
_THUMB_LABEL_PADDING = 0.22
_THUMB_TOP_MIN = 1.2

_SLIDE_WIDTH_IN = SLIDE_WIDTH_EMU / 914400
_SLIDE_HEIGHT_IN = SLIDE_HEIGHT_EMU / 914400


def thumbnail_title_fallback(title: str, style: Style):
    """語の切れ目では 3 行に収まらない題を、確実に収まる大きさまで小さくする。

    ここへ来るのは、題が長すぎる場合だけ。以前はいちばん小さい大きさを選んで
    幅で折り返していたが、**幅の折り返しは最後の 1 行だけ幅を見ない**
    (`text_wrap._hard_wrap` が残りを全部そこへ入れる)。そのため長い題は
    最後の行がスライドの外まで伸びていた。警告も出ないので、画像を見るまで
    分からない ——「サムネイルのレイアウトが崩れる」として出ていたのがこれ。

    ここでは **収まったことを確かめてから返す**。行数・幅・高さの 3 つを見る。
    """
    width_pt = inches_to_pt(_THUMB_WIDTH)
    size = float(_THUMB_TITLE_SIZES[-1])
    while size >= _THUMB_TITLE_HARD_MIN:
        lines = text_wrap.fit_lines(title, size, width_pt, _THUMB_TITLE_MAX_LINES)
        if _thumbnail_lines_fit(lines, size, width_pt, style):
            return size, lines
        size -= _THUMB_TITLE_STEP
    # 3 行に収まらないほど長い題。ここまで来ると読める大きさではないが、
    # 切り詰めれば書いていないものを出すことになり、はみ出させれば崩れる。
    # **崩さないほう** を選び、読めないことは警告で知らせる
    # (`thumbnail_title_is_cramped` -> 「短い題を --title で指定してください」)。
    size = _THUMB_TITLE_HARD_MIN
    return size, text_wrap.fit_lines(title, size, width_pt, _THUMB_TITLE_MAX_LINES)


def _thumbnail_lines_fit(lines, size: float, width_pt: float, style: Style) -> bool:
    """折り返した結果が、行数・幅・高さのすべてで収まっているか。"""
    if len(lines) > _THUMB_TITLE_MAX_LINES:
        return False
    if any(metrics.text_width_em(line) * size > width_pt for line in lines):
        return False
    return len(lines) * style.line_height_pt(size, 1.0) / 72.0 <= _THUMB_TITLE_MAX_HEIGHT


def thumbnail_title_is_cramped(title: str, style: Optional[Style] = None) -> bool:
    """その題が、サムネイルとして読める大きさに収まらないか。

    一覧では小さく表示されるため、ここまで小さくなった題は実際には読めない。
    崩さずに出すことはできるので止めはしないが、知らせる。
    """
    style = style or Style()
    for max_lines in (2, _THUMB_TITLE_MAX_LINES):
        for size in _THUMB_TITLE_SIZES:
            lines = text_wrap.natural_fit_lines(title, size, inches_to_pt(_THUMB_WIDTH), max_lines)
            if lines is not None and len(lines) <= max_lines:
                return False
    return True


def render_deck(deck: Deck, output_path: str, style: Optional[Style] = None) -> str:
    Renderer(style or Style()).render(deck, output_path)
    return output_path


class Renderer:
    def __init__(self, style: Style) -> None:
        self.style = style
        # フッタ用。`render` が呼ばれたときに、デッキの内容で置き換える。
        self._deck_title = ""
        self._total = 1

    def render(self, deck: Deck, output_path: str) -> None:
        prs = Presentation()
        prs.slide_width = Emu(SLIDE_WIDTH_EMU)
        prs.slide_height = Emu(SLIDE_HEIGHT_EMU)
        prs.core_properties.title = deck.title

        # フッタ(資料名とページ番号)のために、全体の枚数と題を持っておく。
        self._deck_title = deck.title
        self._total = len(deck.slides)

        for number, slide in enumerate(deck.slides, start=1):
            self._render_slide(prs, slide, number)

        prs.save(output_path)

    # -- スライド種別ごとの描画 -----------------------------------------
    def _render_slide(self, prs: Presentation, slide: Slide, number: int = 1) -> None:
        if slide.kind == KIND_TITLE:
            pptx_slide = self._title_slide(prs, slide)
        elif slide.kind == KIND_THUMBNAIL:
            pptx_slide = self._thumbnail_slide(prs, slide)
        elif slide.kind == KIND_SECTION:
            pptx_slide = self._section_slide(prs, slide)
            self._draw_footer(pptx_slide, number)
        else:
            pptx_slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_ONLY])
            self._paint_background(pptx_slide, self.style.theme.background)
            self._draw_slide_title(pptx_slide, slide.title)
            self._draw_footer(pptx_slide, number)
            # 1 枚に 1 つでも複数でも、置き場所の決め方は同じ(`layout.fit`)。
            # スライド自身が中身を持つ場合は、1 つだけ並べたものとして扱う。
            parts = slide.parts or [slide]
            placed = layout_mod.fit(parts, self.style, layout_mod.body_box(self.style))
            for part in placed.parts:
                self._draw_content(pptx_slide, part.content, part.box, slide.continued)

        if slide.notes:
            pptx_slide.notes_slide.notes_text_frame.text = slide.notes

    def _title_slide(self, prs: Presentation, slide: Slide):
        """表紙。動画の最初の 1 枚として、地を塗り分けられるようにする。"""
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE])
        style, theme = self.style, self.style.theme
        filled = theme.filled_cover
        self._paint_background(
            pptx_slide, theme.cover_background if filled else theme.background
        )
        title_ph = pptx_slide.shapes.title
        bottom = theme.cover_title_bottom
        top, height = self._cover_title_box(slide.title, style.deck_title_size, bottom)
        _place(title_ph, _COVER_LEFT, top, _COVER_WIDTH, height)
        self._fill_text(
            title_ph.text_frame,
            [[Run(self._wrapped(slide.title, style.deck_title_size))]],
            size=style.deck_title_size,
            color=theme.cover_title if filled else style.color_title,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.BOTTOM,
        )

        subtitle_ph = _placeholder(pptx_slide, 1)
        if slide.subtitle and subtitle_ph is not None:
            _place(subtitle_ph, _COVER_LEFT, bottom + _COVER_SUBTITLE_GAP, _COVER_WIDTH, 0.9)
            self._fill_text(
                subtitle_ph.text_frame,
                [[Run(slide.subtitle)]],
                size=style.deck_subtitle_size,
                color=theme.cover_subtitle if filled else style.color_muted,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.TOP,
            )
        elif subtitle_ph is not None:
            _remove_shape(subtitle_ph)

        self._draw_rule(
            pptx_slide,
            (_COVER_WIDTH + 2 * _COVER_LEFT - _COVER_RULE_WIDTH) / 2,
            bottom + _COVER_RULE_GAP,
            _COVER_RULE_WIDTH,
            style.title_rule_height,
            theme.cover_accent if filled else style.color_accent,
        )
        if filled:
            # 下端の帯。文字だけの画面に、下の重みを付ける。
            self._draw_rule(
                pptx_slide,
                0,
                _SLIDE_HEIGHT_IN - _COVER_BAND_HEIGHT,
                _SLIDE_WIDTH_IN,
                _COVER_BAND_HEIGHT,
                theme.cover_accent,
            )
        return pptx_slide

    def _thumbnail_slide(self, prs: Presentation, slide: Slide):
        """サムネイル(動画の外側で使う 1 枚絵)。

        画面いっぱいに題を置く。一覧では小さく表示されるため、表紙より文字を
        大きくし、左ぞろえにして、行の始まりを追いやすくする。
        """
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BLANK])
        style, theme = self.style, self.style.theme
        filled = theme.filled_cover
        self._paint_background(
            pptx_slide, theme.cover_background if filled else theme.background
        )
        accent = theme.cover_accent if filled else style.color_accent

        size, lines = self._thumbnail_title(slide.title)
        title_height = len(lines) * style.line_height_pt(size, 1.0) / 72.0
        label_height = _THUMB_LABEL_HEIGHT + _THUMB_LABEL_GAP if slide.label else 0.0
        block = label_height + title_height + _THUMB_RULE_GAP + _THUMB_RULE_HEIGHT
        if slide.subtitle:
            block += _THUMB_SUBTITLE_GAP + _THUMB_SUBTITLE_SIZE * 1.4 / 72.0
        # バッジから副題までを 1 かたまりとして、上下の真ん中に置く。
        top = max(_THUMB_TOP_MIN, (_SLIDE_HEIGHT_IN - block) / 2) + label_height

        if slide.label:
            self._draw_thumbnail_label(pptx_slide, slide.label, accent, top)

        box = pptx_slide.shapes.add_textbox(
            Inches(_THUMB_LEFT), Inches(top), Inches(_THUMB_WIDTH), Inches(title_height + 0.2)
        )
        self._fill_text(
            box.text_frame,
            [[Run("\n".join(lines))]],
            size=size,
            color=theme.cover_title if filled else style.color_title,
            bold=True,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
            line_spacing=1.0,
        )

        bottom = top + title_height + _THUMB_RULE_GAP
        self._draw_rule(
            pptx_slide, _THUMB_LEFT, bottom, _THUMB_RULE_WIDTH, _THUMB_RULE_HEIGHT, accent
        )
        if slide.subtitle:
            sub = pptx_slide.shapes.add_textbox(
                Inches(_THUMB_LEFT),
                Inches(bottom + _THUMB_RULE_HEIGHT + _THUMB_SUBTITLE_GAP),
                Inches(_THUMB_WIDTH),
                Inches(0.6),
            )
            self._fill_text(
                sub.text_frame,
                [[Run(slide.subtitle)]],
                size=_THUMB_SUBTITLE_SIZE,
                color=theme.cover_subtitle if filled else style.color_muted,
                align=PP_ALIGN.LEFT,
            )
        if filled:
            self._draw_rule(
                pptx_slide,
                0,
                _SLIDE_HEIGHT_IN - _THUMB_BAND_HEIGHT,
                _SLIDE_WIDTH_IN,
                _THUMB_BAND_HEIGHT,
                accent,
            )
        return pptx_slide

    def _thumbnail_title(self, title: str):
        """サムネイルの題に使う文字の大きさと、折り返した行を選ぶ。

        まず 2 行に収まる中でいちばん大きい文字を探し、無ければ 3 行まで許す。
        行数を増やして文字を大きくするより、2 行で見せるほうが読みやすいため。
        """
        for max_lines in (2, _THUMB_TITLE_MAX_LINES):
            found = self._thumbnail_fit(title, max_lines)
            if found:
                return found
        return thumbnail_title_fallback(title, self.style)

    def _thumbnail_fit(self, title: str, max_lines: int):
        """`max_lines` 行に自然に収まる、いちばん大きい文字を探す(無ければ None)。

        語の途中で切ってまで大きくはしない。文字を小さくすれば自然な位置で
        折り返せるので、そちらを選ぶ。
        """
        for size in _THUMB_TITLE_SIZES:
            lines = text_wrap.natural_fit_lines(
                title, size, inches_to_pt(_THUMB_WIDTH), max_lines
            )
            if lines is None or len(lines) > max_lines:
                continue
            if len(lines) * self.style.line_height_pt(size, 1.0) / 72.0 > _THUMB_TITLE_MAX_HEIGHT:
                continue
            return size, lines
        return None

    def _draw_thumbnail_label(self, pptx_slide, label: str, accent, title_top: float) -> None:
        """左上に置く短い文字(教材名・回数など)。塗った箱に白抜きで入れる。"""
        style = self.style
        width = (
            metrics.text_width_em(label) * _THUMB_LABEL_SIZE / 72.0 + _THUMB_LABEL_PADDING * 2
        )
        height = _THUMB_LABEL_HEIGHT
        top = max(0.7, title_top - _THUMB_LABEL_GAP - height)
        badge = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(_THUMB_LEFT),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        _paint(badge, accent)
        self._fill_text(
            badge.text_frame,
            [[Run(label)]],
            size=_THUMB_LABEL_SIZE,
            color=(0xFF, 0xFF, 0xFF),
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    def _section_slide(self, prs: Presentation, slide: Slide):
        """章扉。本文より少し濃い地にして、章の変わり目が分かるようにする。"""
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_SECTION])
        style, theme = self.style, self.style.theme
        self._paint_background(pptx_slide, theme.section_background)
        title_ph = pptx_slide.shapes.title
        size = style.section_title_size
        lines = self._wrap_lines(slide.title, size, style.title_width)
        height = max(_COVER_TITLE_MIN_HEIGHT, len(lines) * style.line_height_pt(size) / 72.0)
        _place(title_ph, style.title_left, _SECTION_CENTER - height / 2, style.title_width, height)
        self._fill_text(
            title_ph.text_frame,
            [[Run("\n".join(lines))]],
            size=size,
            color=theme.section_title,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        for ph in list(pptx_slide.placeholders):
            if ph != title_ph:
                _remove_shape(ph)
        self._draw_rule(
            pptx_slide,
            (style.title_width + 2 * style.title_left - _COVER_RULE_WIDTH) / 2,
            max(_SECTION_RULE_TOP, _SECTION_CENTER + height / 2 + 0.15),
            _COVER_RULE_WIDTH,
            style.title_rule_height,
            style.color_accent,
        )
        return pptx_slide

    def _draw_slide_title(self, pptx_slide, title: str) -> None:
        style, theme = self.style, self.style.theme
        title_ph = pptx_slide.shapes.title
        if title_ph is None:
            return
        if not title:
            _remove_shape(title_ph)
            return
        _place(title_ph, style.title_left, style.title_top, style.title_width, style.title_height)
        self._fill_text(
            title_ph.text_frame,
            [[Run(title)]],
            size=style.slide_title_size,
            color=style.color_title,
            bold=True,
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
            shrink_to_fit=True,
        )
        if theme.accent_rule_width > 0:
            # 幅いっぱいのうすい罫線に、左端だけ濃い罫線を重ねる。
            self._draw_rule(
                pptx_slide,
                style.title_left,
                style.title_rule_top + (style.title_rule_height - _HAIRLINE) / 2,
                style.title_width,
                _HAIRLINE,
                theme.rule,
            )
            self._draw_rule(
                pptx_slide,
                style.title_left,
                style.title_rule_top,
                theme.accent_rule_width,
                style.title_rule_height,
                style.color_accent,
            )
        else:
            self._draw_rule(
                pptx_slide,
                style.title_left,
                style.title_rule_top,
                style.title_width,
                style.title_rule_height,
                style.color_accent,
            )

    def _draw_footer(self, pptx_slide, number: int) -> None:
        """下端に、資料名(左)とページ番号(右)を小さく置く。"""
        theme = self.style.theme
        if not theme.footer:
            return
        style = self.style
        for text, align in (
            (self._footer_title(), PP_ALIGN.LEFT),
            (f"{number} / {self._total}", PP_ALIGN.RIGHT),
        ):
            if not text:
                continue
            box = pptx_slide.shapes.add_textbox(
                Inches(style.title_left),
                Inches(theme.footer_top),
                Inches(style.title_width),
                Inches(0.3),
            )
            # 飾りであって内容ではないので、ナレーションが読み上げないよう
            # 図形の名前に種類を残す(model.shape_name)。
            box.name = shape_name(SHAPE_FOOTER)
            self._fill_text(
                box.text_frame,
                [[Run(text)]],
                size=theme.footer_size,
                color=theme.footer_color,
                align=align,
            )

    def _footer_title(self) -> str:
        """フッタに出す資料名。長い場合は幅に収まるところまでにする。"""
        title = (self._deck_title or "").strip()
        if not title:
            return ""
        avail = inches_to_pt(self.style.title_width * _FOOTER_TITLE_SHARE) / self.style.theme.footer_size
        if metrics.text_width_em(title) <= avail:
            return title
        clipped = ""
        used = 0.0
        for ch in title:
            if used + metrics.char_width(ch) > avail - 1:
                break
            clipped += ch
            used += metrics.char_width(ch)
        return clipped + "…"

    # -- 折り返しと地の色 -----------------------------------------------
    def _wrap_lines(self, text: str, size: float, width: float = _COVER_WIDTH):
        """大きく映るタイトルを、自然な位置で折り返した行にする。"""
        return text_wrap.fit_lines(text, size, inches_to_pt(width), _COVER_TITLE_MAX_LINES)

    def _wrapped(self, text: str, size: float, width: float = _COVER_WIDTH) -> str:
        return "\n".join(self._wrap_lines(text, size, width))

    def _cover_title_box(self, title: str, size: float, bottom: float):
        """行数ぶんの高さを取り、下端をそろえた箱を返す(top, height)。"""
        lines = self._wrap_lines(title, size)
        height = max(_COVER_TITLE_MIN_HEIGHT, len(lines) * self.style.line_height_pt(size) / 72.0)
        return bottom - height, height

    def _paint_background(self, pptx_slide, rgb) -> None:
        """スライドの地を塗る。白のままなら何も書かない(従来の出力と同じにする)。"""
        if tuple(rgb) == (0xFF, 0xFF, 0xFF):
            return
        fill = pptx_slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb)

    def _draw_rule(self, pptx_slide, left, top, width, height, rgb):
        """細い矩形を 1 つ置く(罫線・帯に使う)。"""
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        _paint(shape, rgb)
        return shape

    # -- 本文 -----------------------------------------------------------
    def _draw_content(self, pptx_slide, content: Content, box: Box, continued: bool) -> None:
        """画面に置く 1 つの中身を、渡された場所に描く。"""
        if content.kind == KIND_CODE:
            self._draw_code(pptx_slide, content, box, continued)
        elif content.kind == KIND_TABLE:
            self._draw_table(pptx_slide, content, box, continued)
        elif content.kind == KIND_IMAGE:
            self._draw_image(pptx_slide, content, box)
        elif content.kind == KIND_DIAGRAM:
            self._draw_diagram(pptx_slide, content, box)
        else:
            self._draw_bullets(pptx_slide, content.bullets, box)

    def _draw_bullets(self, pptx_slide, bullets: List[Bullet], box: Box) -> None:
        if not bullets:
            return
        style = self.style
        shape = pptx_slide.shapes.add_textbox(
            Inches(box.left), Inches(box.top), Inches(box.width), Inches(box.height)
        )
        frame = shape.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.TOP
        _zero_insets(frame)

        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            size = style.body_size(bullet.level)
            self._write_runs(paragraph, bullet.runs, size, style.color_body)
            paragraph.line_spacing = style.line_spacing
            if index > 0:
                paragraph.space_before = Pt(style.space_before_pt)
            self._apply_bullet(paragraph, bullet)

        used = layout_mod.bullets_height(bullets, style, box.width_pt)
        if used > box.height_pt:
            scale = max(0.6, box.height_pt / used)
            oxml_utils.shrink_text(frame, scale, line_reduction=0.1 * (1 - scale))

    def _apply_bullet(self, paragraph, bullet: Bullet) -> None:
        style = self.style
        if bullet.kind == PLAIN:
            # 記号を付けない行。ぶら下げも作らず、左端から書き出す。
            oxml_utils.set_indent(paragraph, Inches(style.indent_per_level * bullet.level), 0)
            oxml_utils.set_no_bullet(paragraph)
            return
        indent = Inches(style.indent_per_level * (bullet.level + 1))
        hanging = Inches(style.indent_per_level)
        oxml_utils.set_indent(paragraph, indent, hanging)
        accent = "%02X%02X%02X" % style.color_accent
        if bullet.kind == NUMBER:
            oxml_utils.set_auto_number(paragraph, start_at=_number_of(bullet))
        elif bullet.kind == QUOTE:
            oxml_utils.set_char_bullet(paragraph, "│", font="Arial", color=accent)
        else:
            char = _BULLET_CHARS.get(bullet.level, "・")
            oxml_utils.set_char_bullet(paragraph, char, font="Arial", color=accent)

    def _draw_code(self, pptx_slide, content: Content, box: Box, continued: bool) -> None:
        style = self.style
        lines = content.code.split("\n")
        height = min(box.height, layout_mod.content_height(content, style, box.width))
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(box.left),
            Inches(box.top),
            Inches(box.width),
            Inches(height),
        )
        # ナレーション生成がコードをそのまま読み上げず、画面の案内文にできるよう、
        # 図形の名前に種類を残す(model.shape_name)。
        shape.name = shape_name(SHAPE_CODE, continued, content.code_lang)
        _paint(shape, style.color_code_bg)
        frame = shape.text_frame
        frame.word_wrap = False
        frame.vertical_anchor = MSO_ANCHOR.TOP
        frame.margin_left = Inches(0.22)
        frame.margin_right = Inches(0.22)
        frame.margin_top = Inches(0.16)
        frame.margin_bottom = Inches(0.16)

        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.line_spacing = style.code_line_spacing
            paragraph.alignment = PP_ALIGN.LEFT
            oxml_utils.set_no_bullet(paragraph)
            run = paragraph.add_run()
            run.text = line if line else " "
            run.font.size = Pt(style.code_size)
            run.font.color.rgb = RGBColor(*style.color_body)
            oxml_utils.set_run_fonts(run, style.font_mono, style.font_mono)

        widest = max((metrics.text_width_em(line) for line in lines), default=0)
        avail_pt = box.width_pt - 40
        if widest * style.code_size > avail_pt:
            scale = max(0.55, avail_pt / (widest * style.code_size))
            oxml_utils.shrink_text(frame, scale)

    def _draw_table(self, pptx_slide, content: Content, box: Box, continued: bool) -> None:
        style = self.style
        header = content.table_header
        rows = content.table_rows
        n_cols = max([len(header)] + [len(r) for r in rows]) or 1
        n_rows = (1 if header else 0) + len(rows)
        if n_rows == 0:
            return
        height = min(box.height, layout_mod.content_height(content, style, box.width))
        shape = pptx_slide.shapes.add_table(
            n_rows,
            n_cols,
            Inches(box.left),
            Inches(box.top),
            Inches(box.width),
            Inches(height),
        )
        shape.name = shape_name(SHAPE_TABLE, continued)
        table = shape.table
        table.first_row = bool(header)

        all_rows = ([header] if header else []) + rows
        for r, row in enumerate(all_rows):
            for c in range(n_cols):
                cell = table.cell(r, c)
                text = row[c] if c < len(row) else ""
                frame = cell.text_frame
                frame.word_wrap = True
                paragraph = frame.paragraphs[0]
                oxml_utils.set_no_bullet(paragraph)
                run = paragraph.add_run()
                run.text = text
                run.font.size = Pt(style.table_size)
                run.font.bold = bool(header) and r == 0
                oxml_utils.set_run_fonts(run, style.font_latin, style.font_ea)
                cell.margin_left = Inches(0.1)
                cell.margin_right = Inches(0.1)

    def _draw_diagram(self, pptx_slide, content: Content, box: Box) -> None:
        """図解を図形として描く。

        文字を等幅で並べた ASCII アートは、日本語と罫線が同じ幅で描かれる
        フォントを前提にしていて、資料 -> PDF -> 画像と変換する間にどこかで
        崩れる。ここでは箱と矢印を図形として置くので、フォントに依存しない。
        """
        items = layout_mod.diagram_items(content)
        if not items:
            return
        style = self.style
        frame_shape = content.diagram_shape == DIAGRAM_FRAME

        # 並べ方も大きさも layout が決める。場所が足りなければ縮み、余っていれば
        # 見出しの大きさまでは広がる。見積りと描画がずれると図が下の中身に重なり、
        # それは画像を見るまで分からないので、数え方を 2 か所に置かない。
        geometry = layout_mod.diagram_geometry(content, style, box.width, box.height)
        item_h = geometry.item_height
        gap = geometry.gap
        pad = geometry.padding
        font_pt = geometry.font_pt

        inner_w = geometry.item_width
        total_w = geometry.width
        total_h = geometry.height
        left = box.left + (box.width - total_w) / 2
        top = box.top + max(0.0, box.height - total_h) / 2

        # 外枠。案内文を二重に出さないよう、図解 1 つにつきこの図形だけに
        # 図解の名前を付ける(narration._is_diagram が見る)。
        outer = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(total_w),
            Inches(total_h),
        )
        # 名前には「どう描いたか」を残す。ナレーションが「上から順に」と
        # 「左から順に」を選ぶのに使う(model.DIAGRAM_FLOW_ACROSS)。
        outer.name = shape_name(
            SHAPE_DIAGRAM,
            language=DIAGRAM_FLOW_ACROSS if geometry.horizontal else content.diagram_shape,
        )
        _zero_insets(outer.text_frame)
        if content.diagram_shape == DIAGRAM_BOUNDARY:
            # 境界図の外枠も、流れ図と同じく目印として置くだけ。図の主役は
            # 真ん中の線なので、外側にもう 1 つ枠があると線が枠の 1 本に見える。
            outer.fill.background()
            _no_outline(outer)
            self._draw_boundary(pptx_slide, content, geometry, left, top)
            return
        if content.diagram_shape == DIAGRAM_LANES:
            # レーン図も同じ。主役は列を分ける縦線で、外枠を描くと 2 本になる。
            outer.fill.background()
            _no_outline(outer)
            self._draw_lanes(pptx_slide, content, geometry, left, top)
            return
        if content.diagram_shape == DIAGRAM_STEPS:
            # 階段図も同じ。主役は段のずれ方で、外枠を描くとずれが枠に
            # 埋もれて「ただ縦に並んだ箱」に見える。
            outer.fill.background()
            _no_outline(outer)
            self._draw_steps(pptx_slide, content, geometry, left, top)
            return
        if frame_shape:
            # 枠図では、この外枠が図そのもの(何が中に入っているかを示す)。
            _paint(outer, style.color_code_bg)
            _outline(outer, style.color_accent, 1.75)
        else:
            # 流れ図では、外枠は案内文を作るときの目印として置くだけなので
            # 描かない(地を塗ると、項目の箱の間にだけ帯が残って figure に見える)。
            outer.fill.background()
            _no_outline(outer)

        item_left = left + pad
        item_top = top + pad
        for index, text in enumerate(items):
            if index > 0:
                if geometry.horizontal:
                    self._draw_diagram_arrow(
                        pptx_slide, item_left, item_top, inner_w, item_h, gap, horizontal=True
                    )
                    item_left += gap
                elif frame_shape:
                    item_top += gap
                else:
                    self._draw_diagram_arrow(
                        pptx_slide, item_left, item_top, inner_w, item_h, gap, horizontal=False
                    )
                    item_top += gap
            self._draw_diagram_item(pptx_slide, text, item_left, item_top, inner_w, item_h, font_pt)
            if geometry.horizontal:
                item_left += inner_w
            else:
                item_top += item_h

    def _draw_boundary(
        self, pptx_slide, content: Content, geometry, left: float, top: float
    ) -> None:
        """境界図を描く。上の箱、線とそれをまたぐもの、下の箱の順に置く。

        線は図全体の幅いっぱいに引く(箱より広い)。線が箱の幅で止まると、
        「箱と箱の区切り」に見えて、越える・越えないの話に見えない。
        またぐものの札は線の上に載せず、**線を挟んでどちら側に置くか** で
        向きを表す(`_draw_boundary_crossings`)。線を 1 本のまま残すため。
        """
        style = self.style
        lines = [text for text in content.diagram_items if text.strip()]
        parts = boundary_parts(lines)
        item_h = geometry.item_height
        gap = geometry.gap
        width = geometry.item_width
        # 箱は線より内側に置く(線だけが左右へはみ出す)。
        box_left = left + (geometry.width - width) / 2

        y = top
        for text in [t for t in parts.upper if t.strip()]:
            self._draw_diagram_item(
                pptx_slide, text, box_left, y, width, item_h, geometry.font_pt
            )
            y += item_h + gap
        band_top = y - gap if parts.upper else y

        self._draw_boundary_rule(pptx_slide, left, band_top, geometry)
        self._draw_boundary_crossings(
            pptx_slide, parts.crossings, box_left, band_top, geometry
        )

        y = band_top + geometry.band
        for text in [t for t in parts.lower if t.strip()]:
            self._draw_diagram_item(
                pptx_slide, text, box_left, y, width, item_h, geometry.font_pt
            )
            y += item_h + gap

    def _draw_boundary_rule(self, pptx_slide, left: float, band_top: float, geometry) -> None:
        """境界の線そのもの。帯の中央に、図の幅いっぱいに引く。"""
        style = self.style
        thickness = max(0.02, style.diagram_boundary_rule * geometry.scale)
        rule = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(band_top + (geometry.band - thickness) / 2),
            Inches(geometry.width),
            Inches(thickness),
        )
        rule.name = shape_name(SHAPE_DIAGRAM_ITEM)
        _paint(rule, style.color_accent)
        _no_outline(rule)

    def _draw_boundary_crossings(
        self, pptx_slide, crossings, left: float, band_top: float, geometry
    ) -> None:
        """線をまたぐもの。矢印を帯いっぱいに立て、その横に札を置く。

        矢印は帯いっぱいに立てて、線の上下へ突き抜けさせる。帯の中で線に
        触れているだけだと、またいでいるようには見えない。

        札は **線を挟んだどちら側に置くかで向きを表す。** 下りるものは線の上に、
        戻るものは線の下に置く。線の上に札を載せると、そこだけ線が途切れて
        見え、**1 本の線** という図の主題が消える(それが図の主役なので)。
        """
        if not crossings:
            return
        style = self.style
        column = geometry.item_width / len(crossings)
        arrow_w = style.diagram_crossing_arrow * geometry.scale
        pad = geometry.gap
        half = geometry.band / 2
        for index, crossing in enumerate(crossings):
            column_left = left + column * index
            shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW if crossing.down else MSO_SHAPE.UP_ARROW,
                Inches(column_left + pad),
                Inches(band_top + geometry.band * 0.06),
                Inches(arrow_w),
                Inches(geometry.band * 0.88),
            )
            shape.name = shape_name(SHAPE_DIAGRAM_ITEM)
            _paint(shape, style.color_accent)
            _no_outline(shape)
            if not crossing.label:
                continue
            label = pptx_slide.shapes.add_textbox(
                Inches(column_left + pad * 2 + arrow_w),
                Inches(band_top if crossing.down else band_top + half),
                Inches(max(0.3, column - pad * 3 - arrow_w)),
                Inches(half),
            )
            # 札の向きは案内文でも使うので、図形の名前に残す(narration が読む)。
            label.name = shape_name(
                SHAPE_DIAGRAM_ITEM, language="down" if crossing.down else "up"
            )
            frame = label.text_frame
            frame.word_wrap = True
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            _zero_insets(frame)
            self._fill_text(
                frame,
                [[Run(crossing.label)]],
                size=geometry.font_pt * style.diagram_crossing_ratio,
                color=style.color_accent,
                align=PP_ALIGN.LEFT,
            )

    # -- レーン図 --------------------------------------------------------
    #: レーン図の図形に付ける名前。narration が案内文を組み立て直すときに読む。
    LANE_HEAD = "lane-head"
    LANE_COLUMNS = ("lane-a", "lane-b")

    def _draw_lanes(
        self, pptx_slide, content: Content, geometry, left: float, top: float
    ) -> None:
        """レーン図を描く。左右 2 列、上から下へ 1 手順 1 行。

        **1 行に置く箱は 1 つだけ** で、反対の列のその行は必ず空く。だから
        レーンをまたぐ矢印も、戻りの横棒も、箱の上を通らない(空いている
        ところだけを通る)。行を詰めて 2 つ置くと、この前提が崩れる。
        """
        parts = lane_parts(content.diagram_items)
        if not parts.steps or len(parts.lanes) < 2:
            return
        style = self.style
        item_w = geometry.item_width
        item_h = geometry.item_height
        gap = geometry.gap
        gutter = geometry.gutter

        cols_left = left + (geometry.ret_band if geometry.ret_left else 0.0)
        col_x = (cols_left, cols_left + item_w + gutter)
        rows_top = top + geometry.header

        def row_top(index: int) -> float:
            return rows_top + index * (item_h + gap)

        # 列を分ける縦線。境界図の線と同じ意味を持つので、同じ太さで引く。
        # レーン名の帯も貫かせる(名前のところで途切れると、線が「箱の区切り」
        # に見えて、担当の分かれ目に見えない)。
        thickness = max(0.02, style.diagram_lane_rule * geometry.scale)
        rule = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cols_left + item_w + (gutter - thickness) / 2),
            Inches(top),
            Inches(thickness),
            Inches(geometry.height),
        )
        rule.name = shape_name(SHAPE_DIAGRAM_ITEM)
        _paint(rule, style.color_accent)
        _no_outline(rule)

        self._draw_lane_headers(pptx_slide, parts, geometry, col_x, top)

        for index, step in enumerate(parts.steps):
            column = parts.lanes.index(step.lane)
            if index > 0:
                self._draw_lane_arrow(
                    pptx_slide, parts, geometry, col_x, index, row_top(index) - gap
                )
            self._draw_diagram_item(
                pptx_slide,
                step.text,
                col_x[column],
                row_top(index),
                item_w,
                item_h,
                geometry.font_pt,
                language=self.LANE_COLUMNS[column],
            )
            # 戻りは、出ていく手順のすぐ後ろに置く。**図形の並び順が
            # シナリオの行の順** になっていないと、資料を読み直して案内文を
            # 組み立てるとき(narration)に、戻りがどの手順から出たのかが
            # 分からなくなる(最後の手順から出たことにされる)。
            for ret in parts.returns:
                if ret.after == index:
                    self._draw_lane_return(
                        pptx_slide, parts, geometry, ret, col_x, left, row_top
                    )

    def _draw_lane_headers(self, pptx_slide, parts, geometry, col_x, top: float) -> None:
        """レーン名の帯。**誰の列か** は、図の中でいちばん先に読まれる。"""
        style = self.style
        height = geometry.header * 0.78
        for column, lane in enumerate(parts.lanes[:2]):
            shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(col_x[column]),
                Inches(top),
                Inches(geometry.item_width),
                Inches(height),
            )
            shape.name = shape_name(SHAPE_DIAGRAM_ITEM, language=self.LANE_HEAD)
            _paint(shape, style.color_accent)
            _no_outline(shape)
            self._fill_text(
                shape.text_frame,
                [[Run(lane)]],
                size=geometry.font_pt * style.diagram_crossing_ratio,
                color=(0xFF, 0xFF, 0xFF),
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    # -- 階段図 ----------------------------------------------------------
    #: 階段図の図形に付ける名前。narration が案内文を組み立て直すときに読む。
    STEP_BADGE = "step-badge"
    STEP_BODY = "step-body"
    STEP_REACH = "step-reach"

    def _draw_steps(
        self, pptx_slide, content: Content, geometry, left: float, top: float
    ) -> None:
        """階段図を描く。段を上から下へ、1 段ずつ右へずらして積む。

        **ずらす幅が図の主題** で、縦に積むだけなら流れ図と同じものになる。
        段と段のあいだには蹴込み(riser)を入れて、離れた箱ではなく
        1 つながりの階段に見えるようにする。

        到達点(`←`)は、右の帯に札を置き、その段の右端まで矢印を引く。
        帯の左端を全部の段でそろえるのは、**どこまで届いたかを縦に見比べる**
        ためで、そろっていないと段の長さの違いに見える。
        """
        parts = step_parts(content.diagram_items)
        if not parts.levels:
            return
        style = self.style
        item_w = geometry.item_width
        item_h = geometry.item_height
        gap = geometry.gap
        offset = geometry.offset
        badge = geometry.badge

        def row(index: int):
            return left + offset * index, top + index * (item_h + gap)

        for index, level in enumerate(parts.levels):
            x, y = row(index)
            if index > 0:
                self._draw_step_riser(pptx_slide, x, y - gap, gap, geometry)
            # 札 -> 中身 の順に置く。narration は図形の並び順で読み直すので、
            # ここの順がそのまま「Lv1: それっぽい指摘を出す」に戻る。
            self._draw_step_badge(pptx_slide, level.name, x, y, badge, item_h, geometry)
            self._draw_diagram_item(
                pptx_slide,
                level.text,
                x + badge,
                y,
                item_w - badge,
                item_h,
                geometry.font_pt,
                language=self.STEP_BODY,
            )
            for reach in parts.reaches:
                if reach.after == index:
                    self._draw_step_reach(pptx_slide, reach, x + item_w, y, left, geometry)

    def _draw_step_badge(
        self, pptx_slide, name: str, x: float, y: float, width: float, height: float, geometry
    ) -> None:
        """段の名前の札。段そのものより先に読まれるので、地を塗って強くする。"""
        style = self.style
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
        )
        shape.name = shape_name(SHAPE_DIAGRAM_ITEM, language=self.STEP_BADGE)
        _paint(shape, style.color_accent)
        _no_outline(shape)
        self._fill_text(
            shape.text_frame,
            [[Run(name)]],
            size=geometry.font_pt,
            color=(0xFF, 0xFF, 0xFF),
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    def _draw_step_riser(self, pptx_slide, x: float, y: float, height: float, geometry) -> None:
        """段と段をつなぐ蹴込み。次の段の左端から、上の段の底まで立てる。

        これが無いと、ずらして置いた箱が **別々に浮いて** 見える。階段として
        読めるかどうかは、段そのものではなくここで決まる。
        """
        style = self.style
        thickness = max(0.02, style.diagram_lane_rule * geometry.scale)
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(thickness),
            Inches(height),
        )
        shape.name = shape_name(SHAPE_DIAGRAM_ITEM)
        _paint(shape, style.color_accent)
        _no_outline(shape)

    def _draw_step_reach(
        self, pptx_slide, reach, right: float, y: float, left: float, geometry
    ) -> None:
        """到達点。右の帯に札を置き、その段の右端まで左向きの矢印を引く。

        矢じりを段の縁まで届かせる(帯の中で止めると、どの段のことか
        決まらない —— レーン図の戻りで直したのと同じところ)。
        """
        style = self.style
        band_left = left + geometry.width - geometry.reach_band
        thickness = min(0.26 * geometry.scale, geometry.item_height * 0.34)
        length = max(0.12, band_left - right)
        arrow = pptx_slide.shapes.add_shape(
            MSO_SHAPE.LEFT_ARROW,
            Inches(right),
            Inches(y + (geometry.item_height - thickness) / 2),
            Inches(length),
            Inches(thickness),
        )
        arrow.name = shape_name(SHAPE_DIAGRAM_ITEM)
        _paint(arrow, style.color_accent)
        _no_outline(arrow)
        if not reach.label:
            return
        label = pptx_slide.shapes.add_textbox(
            Inches(band_left + geometry.reach_band * (1 - style.diagram_step_reach_ratio)),
            Inches(y),
            Inches(geometry.reach_band * style.diagram_step_reach_ratio),
            Inches(geometry.item_height),
        )
        label.name = shape_name(SHAPE_DIAGRAM_ITEM, language=self.STEP_REACH)
        frame = label.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _zero_insets(frame)
        self._fill_text(
            frame,
            [[Run(reach.label)]],
            size=geometry.font_pt * style.diagram_crossing_ratio,
            color=style.color_accent,
            align=PP_ALIGN.LEFT,
        )

    def _draw_lane_arrow(
        self, pptx_slide, parts, geometry, col_x, index: int, band_top: float
    ) -> None:
        """手順から次の手順へ引く矢印。

        同じレーンなら下向き。レーンが変わるところでは、**縦線をまたぐ横向き**
        にする。担当が移ったことは、矢印が線を越えることでしか見えない。
        """
        style = self.style
        before = parts.lanes.index(parts.steps[index - 1].lane)
        after = parts.lanes.index(parts.steps[index].lane)
        if before == after:
            self._draw_diagram_arrow(
                pptx_slide,
                col_x[before],
                band_top,
                geometry.item_width,
                geometry.item_height,
                geometry.gap,
                horizontal=False,
            )
            return
        # 列の中央から中央まで引く。列の縁から縁までだと、線をまたいだことは
        # 見えても、どの箱からどの箱へ移ったのかが見えない。
        thickness = min(0.30 * geometry.scale, geometry.gap * 0.72)
        start = col_x[before] + geometry.item_width / 2
        end = col_x[after] + geometry.item_width / 2
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW if after > before else MSO_SHAPE.LEFT_ARROW,
            Inches(min(start, end)),
            Inches(band_top + (geometry.gap - thickness) / 2),
            Inches(abs(end - start)),
            Inches(thickness),
        )
        shape.name = shape_name(SHAPE_DIAGRAM_ITEM)
        _paint(shape, style.color_accent)
        _no_outline(shape)

    def _draw_lane_return(
        self, pptx_slide, parts, geometry, ret, col_x, left: float, row_top
    ) -> None:
        """戻りの矢印。手順の箱を避けて、外側の帯を通って戻り先の行まで上がる。

        横棒 → 縦棒 → 矢じり の 3 つで組む。縦棒の先に矢じりを付けると、
        矢が上を向いたまま終わり、**どの箱へ戻るのか** が示されない。
        """
        style = self.style
        if not (0 <= ret.after < len(parts.steps)):
            return
        source = parts.steps[ret.after]
        target_lane = parts.lanes[1] if source.lane == parts.lanes[0] else parts.lanes[0]
        target = next((i for i, s in enumerate(parts.steps) if s.lane == target_lane), None)
        if target is None:
            return
        src_col = parts.lanes.index(source.lane)
        tgt_col = parts.lanes.index(target_lane)
        item_w = geometry.item_width
        half = geometry.item_height / 2
        y_from = row_top(ret.after) + half
        y_to = row_top(target) + half
        stem = max(0.02, 0.1 * geometry.scale)
        head = max(0.1, style.diagram_crossing_arrow * geometry.scale)

        if geometry.ret_left:
            stem_x = left + geometry.ret_band * 0.76
            src_edge, tgt_edge = col_x[src_col], col_x[tgt_col]
            label_left, label_width = left, geometry.ret_band * style.diagram_lane_label_ratio
        else:
            stem_x = left + geometry.width - geometry.ret_band * 0.76
            src_edge = col_x[src_col] + item_w
            tgt_edge = col_x[tgt_col] + item_w
            label_left = stem_x + head
            label_width = geometry.ret_band * style.diagram_lane_label_ratio

        def bar(x, y, w, h, kind=MSO_SHAPE.RECTANGLE):
            piece = pptx_slide.shapes.add_shape(
                kind, Inches(x), Inches(y), Inches(w), Inches(h)
            )
            piece.name = shape_name(SHAPE_DIAGRAM_ITEM)
            _paint(piece, style.color_accent)
            _no_outline(piece)

        # 出るところ: 手順の箱から帯まで。反対の列のこの行は空いているので、
        # 横棒は箱の上を通らない(1 行 1 箱だから成り立つ)。
        bar(min(stem_x, src_edge), y_from - stem / 2, abs(src_edge - stem_x), stem)
        # 上がるところ。
        bar(stem_x - stem / 2, y_to, stem, max(stem, y_from - y_to))
        # 戻り先へ入るところ。ここだけ矢じりを付ける。
        bar(
            min(stem_x, tgt_edge),
            y_to - head / 2,
            max(head, abs(tgt_edge - stem_x)),
            head,
            MSO_SHAPE.RIGHT_ARROW if geometry.ret_left else MSO_SHAPE.LEFT_ARROW,
        )
        if not ret.label:
            return
        # 札は帯の中で折り返す。戻る理由は文になりがちで、1 行で置くと
        # 図が横に伸びきる(境界図の札と違い、ここは縦に置き場所がある)。
        label = pptx_slide.shapes.add_textbox(
            Inches(label_left),
            Inches(y_to),
            Inches(max(0.4, label_width)),
            Inches(max(0.3, y_from - y_to)),
        )
        label.name = shape_name(SHAPE_DIAGRAM_ITEM, language="up")
        frame = label.text_frame
        frame.word_wrap = True
        _zero_insets(frame)
        self._fill_text(
            frame,
            [[Run(ret.label)]],
            size=geometry.font_pt * style.diagram_crossing_ratio,
            color=style.color_accent,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    def _draw_diagram_item(
        self, pptx_slide, text: str, left: float, top: float, width: float, height: float,
        font_pt: float, language: str = "",
    ) -> None:
        style = self.style
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        # レーン図では、どちらの列の箱かを名前に残す(narration が案内文を
        # 組み立て直すときに、レーン名と手順を結び直すために読む)。
        shape.name = shape_name(SHAPE_DIAGRAM_ITEM, language=language)
        _paint(shape, (0xFF, 0xFF, 0xFF))
        _outline(shape, style.color_accent, 1.25)
        frame = shape.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _zero_insets(frame)
        self._fill_text(
            frame, [[Run(text)]], size=font_pt, color=style.color_body, align=PP_ALIGN.CENTER
        )

    def _draw_diagram_arrow(
        self,
        pptx_slide,
        left: float,
        top: float,
        item_width: float,
        item_height: float,
        gap: float,
        horizontal: bool,
    ) -> None:
        """項目と項目の間に置く矢印。横に並べた流れでは右向きになる。"""
        style = self.style
        if horizontal:
            thickness = min(0.34, item_height)
            shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(left + gap * 0.12),
                Inches(top + (item_height - thickness) / 2),
                Inches(gap * 0.76),
                Inches(thickness),
            )
        else:
            thickness = min(0.34, gap)
            shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(left + (item_width - thickness) / 2),
                Inches(top + gap * 0.12),
                Inches(thickness),
                Inches(gap * 0.76),
            )
        shape.name = shape_name(SHAPE_DIAGRAM_ITEM)
        _paint(shape, style.color_accent)
        _no_outline(shape)

    def _draw_image(self, pptx_slide, content: Content, box: Box) -> None:
        style = self.style
        caption_space = layout_mod.caption_height() if content.image_alt else 0.0
        max_w = Inches(box.width)
        max_h = Inches(max(0.1, box.height - caption_space))
        picture = pptx_slide.shapes.add_picture(content.image_path, 0, 0, width=max_w)
        if picture.height > max_h:
            ratio = max_h / picture.height
            picture.height = int(picture.height * ratio)
            picture.width = int(picture.width * ratio)
        picture.left = int(Inches(box.left) + (Inches(box.width) - picture.width) / 2)
        # 図が場所を余らせた場合は、上に貼り付けず、渡された場所の中で上下も
        # 真ん中に置く(図だけの画面で、下に大きな空きができるのを避ける)。
        spare = Inches(box.height) - picture.height - Inches(caption_space)
        picture.top = int(Inches(box.top) + max(0, spare) / 2)

        if content.image_alt:
            top = picture.top + picture.height + Inches(0.1)
            caption = pptx_slide.shapes.add_textbox(
                Inches(box.left), top, Inches(box.width), Inches(0.35)
            )
            self._fill_text(
                caption.text_frame,
                [[Run(content.image_alt)]],
                size=style.caption_size,
                color=style.color_muted,
                align=PP_ALIGN.CENTER,
            )

    # -- テキスト共通 ---------------------------------------------------
    def _fill_text(
        self,
        frame,
        paragraphs: List[List[Run]],
        size: float,
        color,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
        shrink_to_fit: bool = False,
        line_spacing: Optional[float] = None,
    ) -> None:
        frame.word_wrap = True
        frame.vertical_anchor = anchor
        _zero_insets(frame)
        for index, runs in enumerate(paragraphs):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = align
            if line_spacing is not None:
                paragraph.line_spacing = line_spacing
            oxml_utils.set_no_bullet(paragraph)
            self._write_runs(paragraph, runs, size, color, bold=bold)
        if shrink_to_fit:
            text = "".join(r.text for runs in paragraphs for r in runs)
            width_pt = metrics.text_width_em(text) * size
            avail = self.style.title_width * 72.0
            if width_pt > avail:
                oxml_utils.shrink_text(frame, max(0.6, avail / width_pt))

    def _write_runs(self, paragraph, runs: List[Run], size: float, color, bold: bool = False) -> None:
        for source in runs:
            # 記事の中の改行は、段落を分けずに行だけを変える(`<a:br/>`)。
            for index, piece in enumerate(source.text.split("\n")):
                if index:
                    oxml_utils.add_line_break(paragraph)
                if piece:
                    self._write_run(paragraph, source.with_text(piece), size, color, bold)

    def _write_run(self, paragraph, source: Run, size: float, color, bold: bool = False) -> None:
        style = self.style
        run = paragraph.add_run()
        run.text = source.text
        font = run.font
        font.size = Pt(size * (0.94 if source.code else 1.0))
        font.bold = bold or source.bold
        font.italic = source.italic
        if source.link:
            run.hyperlink.address = source.link
            font.color.rgb = RGBColor(*style.color_accent)
        else:
            font.color.rgb = RGBColor(*color)
        if source.code:
            oxml_utils.set_run_fonts(run, style.font_mono, style.font_mono)
        else:
            oxml_utils.set_run_fonts(run, style.font_latin, style.font_ea)


# ---------------------------------------------------------------------------
# 小さなヘルパ
# ---------------------------------------------------------------------------


def _placeholder(pptx_slide, idx: int):
    for shape in pptx_slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None


def _remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def _paint(shape, rgb) -> None:
    """図形を単色で塗る。テーマ由来の影や枠線は明示的に打ち消す。"""
    # add_shape が付ける <p:style>(テーマの効果参照)を外さないと、
    # LibreOffice では effectLst の指定より優先されて影が残る。
    style_element = shape._element.find(qn("p:style"))
    if style_element is not None:
        shape._element.remove(style_element)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.fill.background()
    shape.shadow.inherit = False


def _outline(shape, rgb, width_pt: float) -> None:
    line = shape.line
    line.color.rgb = RGBColor(*rgb)
    line.width = Pt(width_pt)


def _no_outline(shape) -> None:
    shape.line.fill.background()


def _zero_insets(frame) -> None:
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0


def _number_of(bullet: Bullet) -> Optional[int]:
    if not bullet.number:
        return None
    digits = bullet.number.rstrip(".)")
    return int(digits) if digits.isdigit() else None


def _place(shape, left: float, top: float, width: float, height: float) -> None:
    shape.left = Inches(left)
    shape.top = Inches(top)
    shape.width = Inches(width)
    shape.height = Inches(height)
