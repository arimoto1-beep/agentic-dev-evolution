"""動画の外側で使う 1 枚絵(YouTube のサムネイル)を作る。

投稿するとき、動画とは別にサムネイルの画像が要る。資料と別の道具で作ると、
題の書き方も色も動画とずれてしまうため、ここでも同じ描画(`renderer`)と同じ
見た目(`style.Theme`)を使い、資料と地続きの 1 枚にする。

    記事 / 教材シナリオ / 資料(.pptx) --> 題・副題 --> thumbnail.png

作り方はスライド画像と同じで、1 枚だけの .pptx を作って PDF 経由で画像にする
(`slide_images`)。フォントの見え方まで資料と同じになる。

大きさの既定は 1280x720(YouTube が推奨する 16:9)。一覧では小さく表示される
ため、題は表紙より大きい文字で、収まる範囲でいちばん大きい大きさが選ばれる。
"""

from __future__ import annotations

import os
import shutil
import tempfile
from dataclasses import dataclass
from typing import Optional

from . import slide_images
from .model import KIND_THUMBNAIL, Deck, Slide
from .renderer import render_deck
from .slide_images import ImageOptions, SlideImageError
from .style import Style

#: YouTube が推奨する大きさ(16:9)。
DEFAULT_WIDTH = 1280

#: 題を取り出せる入力。
SOURCE_SUFFIXES = (".md", ".markdown", ".pptx")


class ThumbnailError(RuntimeError):
    """サムネイルの生成に失敗した場合。"""


@dataclass
class Thumbnail:
    """サムネイルに出す中身。"""

    title: str
    subtitle: str = ""
    label: str = ""

    def to_deck(self) -> Deck:
        return Deck(
            title=self.title,
            slides=[
                Slide(
                    kind=KIND_THUMBNAIL,
                    title=self.title,
                    subtitle=self.subtitle,
                    label=self.label,
                )
            ],
        )


@dataclass
class ThumbnailResult:
    path: str
    width: int
    height: int
    pptx_path: Optional[str] = None  # --keep-pptx を指定した場合のみ残る


def export_thumbnail(
    thumbnail: Thumbnail,
    output_path: str,
    style: Optional[Style] = None,
    width: int = DEFAULT_WIDTH,
    fmt: str = "png",
    soffice_path: Optional[str] = None,
    force: bool = False,
    keep_pptx: bool = False,
    timeout: float = 180,
) -> ThumbnailResult:
    """サムネイルの画像を書き出す。"""
    if not thumbnail.title.strip():
        raise ThumbnailError("サムネイルの題が空です。--title で指定してください。")
    if os.path.exists(output_path) and not force:
        raise ThumbnailError(
            f"出力先が既に存在します: {output_path}\n上書きする場合は --force を付けてください。"
        )

    options = ImageOptions(width=width, fmt=fmt, prefix="thumbnail_", digits=1)
    options.validate()
    outdir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(outdir, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="note2slides-thumb-") as workdir:
        pptx_path = os.path.join(workdir, "thumbnail.pptx")
        render_deck(thumbnail.to_deck(), pptx_path, style=style)
        try:
            result = slide_images.export_slide_images(
                pptx_path,
                os.path.join(workdir, "image"),
                options=options,
                soffice_path=soffice_path,
                force=True,
                timeout=timeout,
            )
        except SlideImageError as exc:
            raise ThumbnailError(str(exc)) from exc
        image = result.images[0]
        shutil.copyfile(image.path, output_path)
        kept = None
        if keep_pptx:
            kept = os.path.splitext(output_path)[0] + ".pptx"
            shutil.copyfile(pptx_path, kept)

    return ThumbnailResult(
        path=output_path, width=image.width, height=image.height, pptx_path=kept
    )


# ---------------------------------------------------------------------------
# 入力から題を取り出す
# ---------------------------------------------------------------------------


def from_source(path: str) -> Thumbnail:
    """記事 / 教材シナリオ / 資料(.pptx)から、題と副題を取り出す。

    サムネイルに出す文字を人が書き直せるよう、ここでは推測を足さない。
    表紙に出ている題と副題を、そのまま持ってくるだけにする。
    """
    suffix = os.path.splitext(path)[1].lower()
    if not os.path.isfile(path):
        raise ThumbnailError(f"入力ファイルが見つかりません: {path}")
    if suffix == ".pptx":
        return _from_pptx(path)
    if suffix in SOURCE_SUFFIXES:
        return _from_markdown(path)
    known = " / ".join(SOURCE_SUFFIXES)
    raise ThumbnailError(f"未対応の入力形式です: {suffix or path}(扱えるのは {known})")


def _from_markdown(path: str) -> Thumbnail:
    from .markdown_reader import parse_article_file
    from .planner import plan_deck
    from .scenario import build_deck, is_scenario_file, read_scenario

    if is_scenario_file(path):
        deck = build_deck(read_scenario(path))
    else:
        deck = plan_deck(parse_article_file(path))
    for slide in deck.slides:
        if slide.title:
            return Thumbnail(title=slide.title, subtitle=slide.subtitle)
    return Thumbnail(title=deck.title)


def _from_pptx(path: str) -> Thumbnail:
    from pptx import Presentation

    from .model import SHAPE_FOOTER, parse_shape_name

    prs = Presentation(path)
    for slide in prs.slides:
        title = slide.shapes.title
        if title is None or not title.text_frame.text.strip():
            continue
        subtitle = ""
        for shape in slide.shapes:
            # python-pptx は同じ図形でも呼ぶたびに別のオブジェクトを返すため、
            # 題そのものかどうかは XML の要素で見る。
            if shape.element is title.element:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            if parse_shape_name(getattr(shape, "name", ""))[0] == SHAPE_FOOTER:
                continue
            text = shape.text_frame.text.strip()
            if text:
                subtitle = text
                break
        return Thumbnail(title=_clean(title.text_frame.text), subtitle=_clean(subtitle))
    raise ThumbnailError(f"題のあるスライドが見つかりませんでした: {path}")


def _clean(text: str) -> str:
    """資料の中で折り返された題を、1 行に戻す。

    表紙の題は読みやすい位置で折り返してあるが(`text_wrap`)、サムネイルでは
    文字の大きさが変わるので、折り返しはやり直す。日本語は行をそのままつなぎ、
    英単語のところだけ空白を戻す(空白で切った改行だったため)。
    """
    lines = [" ".join(part.split()) for part in text.replace("\v", "\n").split("\n")]
    joined = ""
    for line in [line for line in lines if line]:
        if joined and (_is_word_char(joined[-1]) or _is_word_char(line[0])):
            joined += " "
        joined += line
    return joined


def _is_word_char(ch: str) -> bool:
    """英数字のように、続けて書くと語がつながってしまう文字か。"""
    return ch.isascii() and ch.isalnum()
